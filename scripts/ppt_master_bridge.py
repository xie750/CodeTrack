from __future__ import annotations

import argparse
import json
import math
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape


def _load_request(path: Path) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise SystemExit(f"invalid request json: {exc}") from exc
    if not isinstance(value, dict):
        raise SystemExit("request json must be an object")
    return value


def _clean_text(value: Any, limit: int = 220) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text[:limit]


def _bullets(slide: dict[str, Any]) -> list[str]:
    raw = slide.get("bullets", [])
    if not isinstance(raw, list):
        return []
    return [_clean_text(item, 130) for item in raw if _clean_text(item, 130)][:5]


def _write_metadata(path: Path | None, payload: dict[str, Any]) -> None:
    if not path:
        return
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _resolve_ppt_master_paths(request: dict[str, Any]) -> tuple[Path | None, Path | None, Path | None]:
    raw_home = request.get("ppt_master_home") or ""
    raw_skill_dir = request.get("ppt_master_skill_dir") or ""
    home = Path(str(raw_home)).resolve() if str(raw_home).strip() else None
    skill_dir = Path(str(raw_skill_dir)).resolve() if str(raw_skill_dir).strip() else None
    if home and not skill_dir:
        skill_dir = home / "skills" / "ppt-master"
    script = skill_dir / "scripts" / "svg_to_pptx.py" if skill_dir else None
    if not home and skill_dir:
        home = skill_dir.parent.parent
    return home, skill_dir, script


def _svg_text(text: Any, limit: int = 180) -> str:
    return escape(_clean_text(text, limit), {'"': "&quot;", "'": "&apos;"})


def _wrap_text(text: str, limit: int) -> list[str]:
    value = _clean_text(text, 260)
    if not value:
        return []
    if len(value) <= limit:
        return [value]
    chunks: list[str] = []
    current = ""
    for token in re.split(r"(\s+)", value):
        if not token:
            continue
        if len(current) + len(token) <= limit:
            current += token
            continue
        if current.strip():
            chunks.append(current.strip())
        current = token.strip()
    if current.strip():
        chunks.append(current.strip())
    if len(chunks) == 1 and len(chunks[0]) > limit:
        raw = chunks[0]
        chunks = [raw[index : index + limit] for index in range(0, len(raw), limit)]
    return chunks[:3]


def _text_block(x: int, y: int, lines: list[str], *, size: int, color: str, weight: int = 400, line_gap: int = 34) -> str:
    rows = []
    for index, line in enumerate(lines):
        rows.append(
            f'<text x="{x}" y="{y + index * line_gap}" fill="{color}" '
            f'font-family="Microsoft YaHei, Arial, sans-serif" font-size="{size}" '
            f'font-weight="{weight}">{_svg_text(line, 220)}</text>'
        )
    return "\n".join(rows)


def _slide_svg(
    *,
    title: str,
    subtitle: str,
    bullets: list[str],
    notes: str,
    knowledge_point: str,
    index: int,
    total: int,
) -> str:
    palette = [
        ("#1F6FEB", "#14B8A6", "#EAF4FF"),
        ("#0F766E", "#22C55E", "#E8FCF7"),
        ("#7C3AED", "#0EA5E9", "#F3E8FF"),
        ("#2563EB", "#F59E0B", "#EEF4FF"),
    ]
    primary, accent, tint = palette[index % len(palette)]
    safe_title = _wrap_text(title or "自主学习资源", 24)
    safe_subtitle = _wrap_text(subtitle, 52)
    bullet_rows = bullets[:5]
    bullet_svg = []
    for bullet_index, bullet in enumerate(bullet_rows):
        row_y = 232 + bullet_index * 76
        bullet_svg.append(
            f'<rect x="88" y="{row_y - 28}" width="760" height="54" rx="14" fill="#FFFFFF" stroke="#D8E3F0"/>'
            f'<circle cx="116" cy="{row_y - 2}" r="8" fill="{accent}"/>'
            f'{_text_block(142, row_y + 5, _wrap_text(bullet, 52)[:1], size=22, color="#172033", weight=500)}'
        )
    notes_svg = ""
    if notes:
        notes_lines = _wrap_text(f"讲稿提示：{notes}", 46)[:2]
        notes_svg = (
            '<rect x="88" y="618" width="770" height="56" rx="16" fill="#F6FAFF" stroke="#D8E3F0"/>'
            + _text_block(112, 653, notes_lines, size=15, color="#53657E", line_gap=19)
        )
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
<rect width="1280" height="720" fill="#F8FBFF"/>
<rect x="0" y="0" width="1280" height="18" fill="{primary}"/>
<rect x="58" y="64" width="8" height="592" rx="4" fill="{accent}"/>
<rect x="928" y="86" width="246" height="460" rx="28" fill="{tint}" stroke="#D8E3F0"/>
<text x="958" y="136" fill="{primary}" font-family="Microsoft YaHei, Arial, sans-serif" font-size="22" font-weight="700">学习抓手</text>
<text x="958" y="208" fill="#172033" font-family="Microsoft YaHei, Arial, sans-serif" font-size="36" font-weight="700">{_svg_text(knowledge_point, 40)}</text>
<text x="958" y="306" fill="#53657E" font-family="Microsoft YaHei, Arial, sans-serif" font-size="18">先理解规则</text>
<text x="958" y="346" fill="#53657E" font-family="Microsoft YaHei, Arial, sans-serif" font-size="18">再追踪状态</text>
<text x="958" y="386" fill="#53657E" font-family="Microsoft YaHei, Arial, sans-serif" font-size="18">最后测边界</text>
<rect x="948" y="464" width="112" height="38" rx="19" fill="#FFFFFF" stroke="#D8E3F0"/>
<text x="982" y="489" fill="#53657E" font-family="Microsoft YaHei, Arial, sans-serif" font-size="15" font-weight="700">{index + 1:02d}/{total:02d}</text>
{_text_block(88, 122, safe_title, size=42, color="#172033", weight=800, line_gap=50)}
{_text_block(92, 190, safe_subtitle[:1], size=21, color="#53657E", line_gap=28)}
{''.join(bullet_svg)}
{notes_svg}
</svg>
'''


def _write_spec_lock(path: Path) -> None:
    path.write_text(
        """# Execution Lock

## canvas
- viewBox: 0 0 1280 720
- format: PPT 16:9

## communication
- primary_language: zh-Hans

## colors
- bg: #F8FBFF
- bg_secondary: #EAF4FF
- primary: #1F6FEB
- accent: #14B8A6
- secondary_accent: #22C55E
- text: #172033
- text_secondary: #53657E
- border: #D8E3F0

## typography
- font_family: Arial, "Microsoft YaHei", "PingFang SC", sans-serif
- title_family: Arial, "Microsoft YaHei", "PingFang SC", sans-serif
- body_family: Arial, "Microsoft YaHei", "PingFang SC", sans-serif
- emphasis_family: Arial, "Microsoft YaHei", "PingFang SC", sans-serif
- code_family: Consolas, "Courier New", monospace
- body: 20
- title: 40
- subtitle: 22
- annotation: 14
- footnote: 11
""",
        encoding="utf-8",
    )


def _write_ppt_master_project(request: dict[str, Any], output_path: Path) -> Path:
    workspace_dir = Path(str(request.get("workspace_dir") or output_path.parent)).resolve()
    resource_id = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(request.get("resource_id") or output_path.stem))
    project_dir = workspace_dir / f"ppt_master_project_{resource_id}"
    if project_dir.exists():
        shutil.rmtree(project_dir)
    svg_dir = project_dir / "svg_output"
    notes_dir = project_dir / "notes"
    svg_dir.mkdir(parents=True, exist_ok=True)
    notes_dir.mkdir(parents=True, exist_ok=True)
    _write_spec_lock(project_dir / "spec_lock.md")

    title = _clean_text(request.get("title"), 120) or "自主学习资源"
    message = _clean_text(request.get("message"), 180)
    knowledge_point = _clean_text(request.get("knowledge_point"), 80) or "自主学习"
    slides = request.get("slides", [])
    slides = slides if isinstance(slides, list) and slides else [{"title": title, "subtitle": message, "bullets": [message]}]
    slides = slides[:10]
    for index, raw_slide in enumerate(slides):
        item = raw_slide if isinstance(raw_slide, dict) else {}
        slide_title = _clean_text(item.get("title"), 120) or (title if index == 0 else f"第 {index + 1} 页")
        subtitle = _clean_text(item.get("subtitle"), 150) or (message if index == 0 else "")
        bullets = _bullets(item) or ([message] if message else [])
        notes = _clean_text(item.get("speaker_notes"), 360)
        stem = f"{index + 1:02d}_{re.sub(r'[^A-Za-z0-9]+', '_', slide_title)[:24] or 'slide'}"
        (svg_dir / f"{stem}.svg").write_text(
            _slide_svg(
                title=slide_title,
                subtitle=subtitle,
                bullets=bullets,
                notes=notes,
                knowledge_point=knowledge_point,
                index=index,
                total=len(slides),
            ),
            encoding="utf-8",
        )
        if notes:
            (notes_dir / f"{stem}.md").write_text(notes, encoding="utf-8")

    (project_dir / "metadata.json").write_text(
        json.dumps(
            {
                "title": title,
                "subject": knowledge_point,
                "category": "CodeTrack 自主学习资源",
                "keywords": ["AI专业", "自主学习", knowledge_point],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return project_dir


def _try_render_with_ppt_master(request: dict[str, Any], output_path: Path) -> dict[str, Any] | None:
    ppt_master_home, skill_dir, script = _resolve_ppt_master_paths(request)
    if not script or not script.exists():
        return None
    project_dir = _write_ppt_master_project(request, output_path)
    command = [
        sys.executable,
        str(script),
        str(project_dir),
        "--output",
        str(output_path),
        "--source",
        "svg_output",
        "--pptx-structure",
        "flat",
        "--no-notes",
        "--quiet",
    ]
    result = subprocess.run(
        command,
        cwd=str(skill_dir),
        text=True,
        encoding="utf-8",
        errors="replace",
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=240,
    )
    if result.returncode != 0:
        raise RuntimeError((result.stderr or result.stdout or "").strip()[-1600:])
    if not output_path.exists() or output_path.stat().st_size < 100:
        raise RuntimeError("PPT Master converter did not write a valid PPTX")
    return {
        "project_id": project_dir.name,
        "project_path": str(project_dir),
        "export_path": str(output_path),
        "implementation": "ppt_master_svg_to_pptx",
        "ppt_master_home": str(ppt_master_home or ""),
        "ppt_master_skill_dir": str(skill_dir or ""),
        "stdout_tail": (result.stdout or "").strip()[-800:],
        "stderr_tail": (result.stderr or "").strip()[-800:],
    }


def _add_text(slide, left, top, width, height, text: str, size, color, *, bold=False, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = font
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _shape(slide, kind, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _render_pptx(request: dict[str, Any], output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        raise SystemExit(f"python-pptx is required: {exc}") from exc

    title = _clean_text(request.get("title"), 120) or "自主学习资源"
    message = _clean_text(request.get("message"), 180)
    knowledge_point = _clean_text(request.get("knowledge_point"), 80) or "自主学习"
    slides = request.get("slides", [])
    slides = slides if isinstance(slides, list) else []
    citations = request.get("citations", [])
    citations = citations if isinstance(citations, list) else []

    dark = RGBColor(20, 30, 48)
    blue = RGBColor(31, 111, 235)
    teal = RGBColor(20, 184, 166)
    green = RGBColor(34, 197, 94)
    slate = RGBColor(86, 101, 124)
    pale_blue = RGBColor(236, 246, 255)
    pale_teal = RGBColor(232, 252, 247)
    white = RGBColor(255, 255, 255)
    border = RGBColor(218, 228, 242)
    palette = [blue, teal, green, RGBColor(124, 58, 237), RGBColor(14, 165, 233)]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, raw_slide in enumerate(slides[:10] or [{"title": title, "bullets": [message]}]):
        item = raw_slide if isinstance(raw_slide, dict) else {}
        accent = palette[index % len(palette)]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 251, 255)

        _shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16), accent)
        _shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(0.55), Inches(0.08), Inches(6.25), accent)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.72), Inches(0.52), Inches(1.85), Inches(0.38), pale_blue, border)
        _add_text(slide, Inches(10.86), Inches(0.59), Inches(1.52), Inches(0.2), f"{index + 1:02d} / {max(len(slides), 1):02d}", Pt(9), slate, bold=True)

        if index == 0:
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.06), Inches(3.1), Inches(0.42), pale_teal, border)
            _add_text(slide, Inches(1.14), Inches(1.16), Inches(2.62), Inches(0.18), knowledge_point, Pt(10), teal, bold=True)
            _add_text(slide, Inches(0.95), Inches(1.86), Inches(8.4), Inches(1.16), _clean_text(item.get("title"), 120) or title, Pt(38), dark, bold=True)
            subtitle = _clean_text(item.get("subtitle"), 140) or message or "AI 专业自主学习课件"
            _add_text(slide, Inches(1.0), Inches(3.16), Inches(8.6), Inches(0.46), subtitle, Pt(17), slate)
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.86), Inches(1.78), Inches(2.42), Inches(2.42), accent)
            _add_text(slide, Inches(10.3), Inches(2.5), Inches(1.52), Inches(0.52), "AI\nCOURSE", Pt(24), white, bold=True)
            top = 4.28
            for bullet in _bullets(item)[:3]:
                _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(top), Inches(0.22), Inches(0.22), accent)
                _add_text(slide, Inches(1.38), Inches(top - 0.03), Inches(7.6), Inches(0.32), bullet, Pt(14), dark)
                top += 0.48
        else:
            _add_text(slide, Inches(0.95), Inches(0.7), Inches(8.95), Inches(0.62), _clean_text(item.get("title"), 110) or f"第 {index + 1} 页", Pt(29), dark, bold=True)
            subtitle = _clean_text(item.get("subtitle"), 140)
            if subtitle:
                _add_text(slide, Inches(0.98), Inches(1.31), Inches(8.2), Inches(0.32), subtitle, Pt(13), slate)

            bullets = _bullets(item)
            cols = 2 if len(bullets) >= 4 else 1
            card_width = 5.35 if cols == 2 else 8.5
            for bullet_index, bullet in enumerate(bullets):
                col = bullet_index % cols
                row = math.floor(bullet_index / cols)
                left = Inches(0.98 + col * 5.68)
                top = Inches(1.95 + row * 1.18)
                _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(card_width), Inches(0.88), white, border)
                _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.18), top + Inches(0.22), Inches(0.32), Inches(0.32), accent)
                _add_text(slide, left + Inches(0.68), top + Inches(0.18), Inches(card_width - 0.9), Inches(0.46), bullet, Pt(14), dark)

            notes = _clean_text(item.get("speaker_notes"), 360)
            if notes:
                _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.98), Inches(5.72), Inches(8.7), Inches(0.7), pale_blue, border)
                _add_text(slide, Inches(1.2), Inches(5.87), Inches(8.25), Inches(0.32), f"讲稿提示：{notes}", Pt(10), slate)

            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.18), Inches(1.34), Inches(2.22), Inches(4.72), pale_teal, border)
            _add_text(slide, Inches(10.48), Inches(1.72), Inches(1.58), Inches(0.3), "学习抓手", Pt(13), teal, bold=True)
            _add_text(slide, Inches(10.48), Inches(2.18), Inches(1.58), Inches(1.5), knowledge_point, Pt(24), dark, bold=True)
            _add_text(slide, Inches(10.48), Inches(4.22), Inches(1.55), Inches(0.76), "先理解规则\n再追踪状态\n最后测边界", Pt(11), slate)

        citation_titles = [_clean_text(citation.get("title"), 56) for citation in citations[:2] if isinstance(citation, dict)]
        if citation_titles:
            _add_text(slide, Inches(0.96), Inches(6.74), Inches(9.2), Inches(0.24), "引用：" + "；".join(citation_titles), Pt(8), slate)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="CodeTrack PPT Master bridge")
    parser.add_argument("--request-json", required=True)
    parser.add_argument("--output-pptx", required=True)
    args = parser.parse_args()

    request_path = Path(args.request_json).resolve()
    output_path = Path(args.output_pptx).resolve()
    request = _load_request(request_path)
    request.setdefault("workspace_dir", str(request_path.parent))
    provider_metadata: dict[str, Any]
    try:
        official_metadata = _try_render_with_ppt_master(request, output_path)
    except Exception as exc:
        _render_pptx(request, output_path)
        provider_metadata = {
            "project_id": f"ppt_master_bridge_{request.get('resource_id', 'resource')}",
            "export_path": str(output_path),
            "implementation": "codetrack_ppt_master_bridge_fallback",
            "ppt_master_home": str(request.get("ppt_master_home") or ""),
            "official_converter_error": str(exc)[:800],
        }
    else:
        if official_metadata:
            provider_metadata = official_metadata
        else:
            _render_pptx(request, output_path)
            provider_metadata = {
                "project_id": f"ppt_master_bridge_{request.get('resource_id', 'resource')}",
                "export_path": str(output_path),
                "implementation": "codetrack_ppt_master_bridge_fallback",
                "ppt_master_home": str(request.get("ppt_master_home") or ""),
                "official_converter_error": "PPT Master home or svg_to_pptx.py was not found.",
            }

    metadata_path = Path(str(request.get("metadata_json") or "")).resolve() if request.get("metadata_json") else None
    _write_metadata(metadata_path, provider_metadata)
    print(f"wrote {output_path} via {provider_metadata.get('implementation')}")


if __name__ == "__main__":
    main()
