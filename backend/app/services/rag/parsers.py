from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
import re
import tempfile

from backend.app.services.rag.utils import normalize_text


@dataclass
class ParsedElement:
    element_type: str
    text: str
    page_no: int | None = None
    slide_no: int | None = None
    heading_level: int | None = None
    heading_path: list[str] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)


@dataclass
class ParseResult:
    parser_name: str
    parser_version: str
    elements: list[ParsedElement]


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}


def parse_document(filename: str, content: bytes) -> ParseResult:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"unsupported file type: {extension}")
    if extension in {".txt", ".md", ".markdown"}:
        return parse_text_document(filename, content)
    if extension == ".pdf":
        return parse_pdf_document(content)
    if extension == ".docx":
        return parse_docx_document(content)
    if extension == ".pptx":
        return parse_pptx_document(content)
    raise ValueError(f"unsupported file type: {extension}")


def parse_text_document(filename: str, content: bytes) -> ParseResult:
    text = content.decode("utf-8-sig")
    extension = Path(filename).suffix.lower()
    is_markdown = extension in {".md", ".markdown"}
    heading_path: list[str] = []
    elements: list[ParsedElement] = []
    paragraph_blocks: list[str] = []
    list_blocks: list[str] = []
    table_blocks: list[str] = []
    code_blocks: list[str] = []
    in_code_block = False
    code_fence = ""

    def flush_paragraph() -> None:
        nonlocal paragraph_blocks
        merged = normalize_text("\n".join(paragraph_blocks))
        paragraph_blocks = []
        if merged:
            elements.append(ParsedElement("paragraph", merged, heading_path=list(heading_path)))

    def flush_list() -> None:
        nonlocal list_blocks
        merged = normalize_text("\n".join(list_blocks))
        list_blocks = []
        if merged:
            elements.append(ParsedElement("list", merged, heading_path=list(heading_path)))

    def flush_table() -> None:
        nonlocal table_blocks
        merged = normalize_text("\n".join(table_blocks))
        table_blocks = []
        if merged:
            elements.append(ParsedElement("table", merged, heading_path=list(heading_path), metadata={"format": "markdown"}))

    def flush_code() -> None:
        nonlocal code_blocks, code_fence
        merged = "\n".join(code_blocks).strip()
        metadata = {"format": "markdown_fence"}
        if code_fence:
            metadata["fence"] = code_fence
        code_blocks = []
        code_fence = ""
        if merged:
            elements.append(ParsedElement("code", merged, heading_path=list(heading_path), metadata=metadata))

    def flush_flow() -> None:
        flush_paragraph()
        flush_list()
        flush_table()

    for line in text.splitlines():
        stripped = line.strip()
        if is_markdown and stripped.startswith("```"):
            if in_code_block:
                code_blocks.append(line)
                flush_code()
                in_code_block = False
            else:
                flush_flow()
                in_code_block = True
                code_fence = stripped
                code_blocks.append(line)
            continue
        if in_code_block:
            code_blocks.append(line)
            continue
        if is_markdown:
            match = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
            if match:
                flush_flow()
                level = len(match.group(1))
                title = normalize_text(match.group(2))
                if _looks_like_attribute_heading(title):
                    paragraph_blocks.append(title)
                    continue
                heading_path = heading_path[: level - 1] + [title]
                elements.append(
                    ParsedElement(
                        "heading",
                        title,
                        heading_level=level,
                        heading_path=list(heading_path),
                    )
                )
                continue
        outline_heading = _match_markdown_outline_heading(stripped)
        if outline_heading:
            flush_flow()
            level, title = outline_heading
            if _looks_like_attribute_heading(title):
                paragraph_blocks.append(title)
                continue
            heading_path = heading_path[: level - 1] + [title]
            elements.append(
                ParsedElement(
                    "heading",
                    title,
                    heading_level=level,
                    heading_path=list(heading_path),
                    metadata={"format": "text_outline" if not is_markdown else "markdown_outline"},
                )
            )
            continue
        if is_markdown:
            if _is_markdown_table_line(stripped):
                flush_paragraph()
                flush_list()
                table_blocks.append(line)
                continue
            if table_blocks:
                flush_table()
            if re.match(r"^([-*+]|\d+[.)])\s+", stripped):
                flush_paragraph()
                list_blocks.append(line)
                continue
            if list_blocks:
                flush_list()
        if not stripped:
            flush_flow()
            continue
        paragraph_blocks.append(line)
    if in_code_block:
        flush_code()
    flush_flow()
    return ParseResult(
        parser_name="markdown" if is_markdown else "plain_text",
        parser_version="1",
        elements=elements,
    )


def _is_markdown_table_line(line: str) -> bool:
    if "|" not in line:
        return False
    cells = [cell.strip() for cell in line.strip("|").split("|")]
    if len(cells) < 2:
        return False
    return any(cells) or bool(re.match(r"^[:\-\s|]+$", line))


def _match_markdown_outline_heading(line: str) -> tuple[int, str] | None:
    if not line:
        return None
    patterns = [
        r"^(?P<num>\d+(?:\.\d+)*[.、)]?)\s+\*\*(?P<title>[^*]{2,80})\*\*\s*$",
        r"^\*\*(?P<title>[^*]{2,80})[:：]?\*\*\s*$",
        r"^(?P<num>\d+(?:\.\d+)*[.、)]?)\s+(?P<title>[^。！？!?；;：:]{2,60})[:：]?\s*$",
    ]
    for pattern in patterns:
        match = re.match(pattern, line)
        if not match:
            continue
        title = normalize_text(match.group("title"))
        if not title or _looks_like_list_body(title):
            return None
        number = match.groupdict().get("num") or ""
        normalized_number = re.sub(r"[.、)]$", "", number)
        level = min(3, normalized_number.count(".") + 1) if normalized_number else 2
        return level, title
    return None


def _looks_like_list_body(text: str) -> bool:
    lowered = text.lower()
    if len(text) > 60:
        return True
    if any(mark in lowered for mark in ["http://", "https://", "return ", "def ", "#include"]):
        return True
    return False


def _looks_like_attribute_heading(text: str) -> bool:
    normalized = normalize_text(text).strip(" #>*-　")
    normalized = re.sub(r"^\*\*(.+?)\*\*", r"\1", normalized).strip()
    labels = [
        "常见定位",
        "定位",
        "角色定位",
        "英雄定位",
        "技能特点",
        "玩法特点",
        "打法特点",
        "适合位置",
        "核心特点",
        "特点",
        "简介",
        "说明",
        "备注",
        "注",
    ]
    return any(normalized.startswith(label) for label in labels)


def parse_pdf_document(content: bytes) -> ParseResult:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF parsing") from exc

    elements: list[ParsedElement] = []
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as temp:
        temp.write(content)
        temp.flush()
        doc = fitz.open(temp.name)
        for index, page in enumerate(doc, start=1):
            text = normalize_text(page.get_text("text"))
            if text:
                elements.append(ParsedElement("paragraph", text, page_no=index, metadata={"page": index}))
        doc.close()
    return ParseResult("pymupdf", "1", elements)


def parse_docx_document(content: bytes) -> ParseResult:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for DOCX parsing") from exc

    elements: list[ParsedElement] = []
    heading_path: list[str] = []
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=True) as temp:
        temp.write(content)
        temp.flush()
        doc = Document(temp.name)
        for paragraph in doc.paragraphs:
            text = normalize_text(paragraph.text)
            if not text:
                continue
            style_name = (paragraph.style.name if paragraph.style is not None else "").lower()
            if style_name.startswith("heading"):
                match = re.search(r"(\d+)", style_name)
                level = int(match.group(1)) if match else 1
                heading_path = heading_path[: level - 1] + [text]
                elements.append(ParsedElement("heading", text, heading_level=level, heading_path=list(heading_path)))
            else:
                elements.append(ParsedElement("paragraph", text, heading_path=list(heading_path)))
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [normalize_text(cell.text) for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            table_text = normalize_text("\n".join(rows))
            if table_text:
                elements.append(ParsedElement("table", table_text, heading_path=list(heading_path), metadata={"format": "docx_table"}))
    return ParseResult("python-docx", "1", elements)


def parse_pptx_document(content: bytes) -> ParseResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX parsing") from exc

    elements: list[ParsedElement] = []
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=True) as temp:
        temp.write(content)
        temp.flush()
        deck = Presentation(temp.name)
        for slide_index, slide in enumerate(deck.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                text = normalize_text(getattr(shape, "text", "") or "")
                if text:
                    slide_texts.append(text)
            if not slide_texts:
                continue
            title = slide_texts[0]
            elements.append(
                ParsedElement(
                    "heading",
                    title,
                    slide_no=slide_index,
                    heading_level=1,
                    heading_path=[title],
                    metadata={"slide": slide_index},
                )
            )
            for text in slide_texts[1:]:
                elements.append(
                    ParsedElement(
                        "paragraph",
                        text,
                        slide_no=slide_index,
                        heading_path=[title],
                        metadata={"slide": slide_index},
                    )
                )
    return ParseResult("python-pptx", "1", elements)
