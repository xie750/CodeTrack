import json
import os
import re
import shutil
import subprocess
from pathlib import Path
from typing import Any, TypedDict
from urllib.parse import urljoin
from uuid import uuid4

from sqlalchemy import select
from sqlalchemy.orm import Session

from backend.app.ai.errors import LLMError
from backend.app.ai.llm_client import chat_json
from backend.app.ai.run_recorder import finish_run, new_run_id, record_step, start_run
from backend.app.ai.schemas import AgentRunContext
from backend.app.core.api_response import ApiError
from backend.app.core.config import get_settings
from backend.app.models import (
    AgentRun,
    Course,
    KnowledgeSource,
    LearnerEvent,
    StudentGeneratedResource,
    User,
)
from backend.app.models.entities import utc_now
from backend.app.services.learner_profile import serialize_learner_profile
from backend.app.services.presenton_client import (
    PresentonError,
    fetch_presenton_slides_sync,
    generate_presenton_pptx,
    presenton_configured,
)
from backend.app.services.ppt_master_client import (
    PptMasterError,
    generate_ppt_master_pptx,
    ppt_master_configured,
)
from backend.app.services.submissions import iso

try:
    from langgraph.graph import END, StateGraph
except Exception:  # pragma: no cover - optional dependency fallback
    END = None
    StateGraph = None


WORKFLOW_TYPE = "student_ppt_resource_generation"
GENERIC_WORKFLOW_TYPE = "student_resource_generation"
MIND_MAP_WORKFLOW_TYPE = "student_mind_map_resource_generation"
PROMPT_VERSION = "student_ppt_resource_v0.1"
GENERIC_PROMPT_VERSION = "student_resource_v0.2"
MIND_MAP_PROMPT_VERSION = "student_mind_map_resource_v0.1"
MAX_SOURCE_COUNT = 5
MAX_SOURCE_CHARS = 900

SUPPORTED_RESOURCE_TYPES = {
    "PPT",
    "DOCUMENT",
    "MIND_MAP",
    "PRACTICE_SET",
    "KNOWLEDGE_CARD",
    "PODCAST_SCRIPT",
}

RESOURCE_TYPE_LABELS = {
    "PPT": "PPT",
    "DOCUMENT": "文档",
    "MIND_MAP": "思维导图",
    "PRACTICE_SET": "练习题",
    "KNOWLEDGE_CARD": "知识卡片",
    "PODCAST_SCRIPT": "播客稿",
}


class PptResourceState(TypedDict, total=False):
    run_id: str
    student_id: str
    class_id: str
    course_id: str
    session_id: str | None
    message: str
    knowledge_point: str
    profile: dict[str, Any] | None
    sources: list[KnowledgeSource]
    citations: list[dict[str, Any]]
    slides: list[dict[str, Any]]
    presenton_slides: list[dict[str, Any]]
    title: str
    summary: str
    confidence: float
    file_path: str
    model_content_fallback_error: str
    resource: StudentGeneratedResource
    run: AgentRun


class GenericResourceState(TypedDict, total=False):
    run_id: str
    student_id: str
    class_id: str
    course_id: str
    session_id: str | None
    message: str
    resource_type: str
    knowledge_point: str
    profile: dict[str, Any] | None
    sources: list[KnowledgeSource]
    citations: list[dict[str, Any]]
    title: str
    summary: str
    confidence: float
    render_payload: dict[str, Any]
    file_path: str
    file_format: str
    item_count: int
    resource: StudentGeneratedResource
    run: AgentRun


class MindMapResourceState(TypedDict, total=False):
    run_id: str
    student_id: str
    class_id: str
    course_id: str
    session_id: str | None
    message: str
    resource_type: str
    knowledge_point: str
    profile: dict[str, Any] | None
    profile_focus: dict[str, list[str]]
    sources: list[KnowledgeSource]
    citations: list[dict[str, Any]]
    plan: dict[str, Any]
    render_payload: dict[str, Any]
    title: str
    summary: str
    confidence: float
    item_count: int
    file_path: str
    file_format: str
    risk_flags: list[str]
    resource: StudentGeneratedResource
    run: AgentRun


def _json_dumps(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"), default=str)


def _json_loads(raw: str | None, fallback: Any) -> Any:
    if not raw:
        return fallback
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return fallback


def _new_id(prefix: str) -> str:
    return f"{prefix}_{uuid4().hex[:12]}"


def _trim(text: str | None, limit: int) -> str:
    value = (text or "").strip()
    if len(value) <= limit:
        return value
    return f"{value[:limit]}..."


def _safe_json_list(raw: str | None) -> list[str]:
    value = _json_loads(raw, [])
    return [str(item) for item in value] if isinstance(value, list) else []


def _terms(text: str) -> list[str]:
    words = re.findall(r"[A-Za-z0-9_]{2,}|[\u4e00-\u9fff]{2,}", text)
    grams: list[str] = []
    seen: set[str] = set()
    for word in words:
        lowered = word.lower()
        candidates = [lowered]
        if re.fullmatch(r"[\u4e00-\u9fff]{3,}", word):
            candidates.extend(word[index : index + 2] for index in range(len(word) - 1))
        for item in candidates:
            if item and item not in seen:
                seen.add(item)
                grams.append(item)
    return grams


def _relevance(query: str, source: KnowledgeSource) -> float:
    content = "\n".join(
        [
            source.title or "",
            source.summary or "",
            source.content or "",
            " ".join(_safe_json_list(source.knowledge_points)),
        ]
    ).lower()
    terms = _terms(query)
    if not terms:
        return 0
    return sum(1 for term in terms if term.lower() in content) / max(len(terms), 1)


def _cleanup_topic(raw: str) -> str:
    topic = re.sub(r"\s+", " ", raw).strip(" ：:，,。.；;、")
    topic = re.sub(r"^(帮我|请|生成|整理|制作|做一个|做一份|输出|围绕|关于)+", "", topic, flags=re.IGNORECASE)
    topic = topic.strip(" ：:，,。.；;、")
    suffix_pattern = r"(相关|有关)?的?(PPTX?|pptx?|幻灯片|演示文稿|思维导图|导图|学习地图|学习文档|讲解文档|复习文档|文档|教学讲解|学习资源|复习资料|讲解|资料|资源)$"
    while True:
        cleaned = re.sub(suffix_pattern, "", topic, flags=re.IGNORECASE).strip(" ：:，,。.；;、")
        cleaned = re.sub(r"(相关|有关)的?$", "", cleaned, flags=re.IGNORECASE).strip(" ：:，,。.；;、")
        if cleaned == topic:
            break
        topic = cleaned
    topic = re.sub(r"\s+", " ", topic).strip(" ：:，,。.；;、")
    return topic[:40]


def _extract_requested_topic(message: str) -> str | None:
    text = message.strip()
    patterns = [
        r"(?:把)(.+?)(?:整理|制作|生成|做成|转成|变成)(?:成|为)?(?:.+?)(?:思维导图|导图|学习地图)",
        r"(?:关于|围绕)(.+?)(?:的)?(?:思维导图|导图|学习地图|PPT|ppt|幻灯片|演示文稿|学习文档|讲解文档|复习文档|文档|讲解|资料|资源|$)",
        r"(?:生成|制作|整理|做一份|做一个)(.+?)(?:相关|有关|的)?(?:思维导图|导图|学习地图|PPT|ppt|幻灯片|演示文稿|学习文档|讲解文档|复习文档|文档|资料|资源|$)",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if not match:
            continue
        topic = _cleanup_topic(match.group(1))
        if topic:
            return topic

    topic = _cleanup_topic(text)
    if topic and len(topic) <= 24:
        return topic
    return None


def _normalize_topic_text(text: str) -> str:
    return "".join(re.findall(r"[A-Za-z0-9]+|[\u4e00-\u9fff]+", text.lower()))


def _assert_slides_match_topic(title: str, summary: str, slides: list[dict[str, Any]], knowledge_point: str) -> None:
    topic = _normalize_topic_text(knowledge_point)
    if not topic or knowledge_point == "自主学习主题":
        return
    content = _normalize_topic_text(
        " ".join(
            [
                title,
                summary,
                *[
                    " ".join(
                        [
                            str(slide.get("title", "")),
                            str(slide.get("subtitle", "")),
                            " ".join(str(item) for item in slide.get("bullets", [])),
                        ]
                    )
                    for slide in slides
                ],
            ]
        )
    )
    if topic not in content:
        raise ValueError(f"model slides do not match requested topic: {knowledge_point}")


def _citation(source: KnowledgeSource) -> dict[str, Any]:
    return {
        "source_id": source.id,
        "title": source.title,
        "summary": source.summary,
        "source_type": source.source_type,
        "version": source.version,
        "authority_level": source.authority_level,
    }


def _source_payload(source: KnowledgeSource) -> dict[str, Any]:
    return {
        **_citation(source),
        "chapter": source.chapter,
        "knowledge_points": _safe_json_list(source.knowledge_points),
        "content": _trim(source.content or source.summary, MAX_SOURCE_CHARS),
    }


def _load_sources(db: Session, course_id: str, message: str) -> list[KnowledgeSource]:
    all_sources = list(
        db.scalars(
            select(KnowledgeSource)
            .where(
                KnowledgeSource.course_id == course_id,
                KnowledgeSource.status == "ACTIVE",
                KnowledgeSource.ai_retrievable.is_(True),
                KnowledgeSource.student_visible.is_(True),
            )
            .order_by(KnowledgeSource.authority_level.desc(), KnowledgeSource.updated_at.desc())
        ).all()
    )
    scored = [(source, _relevance(message, source)) for source in all_sources]
    matched = [source for source, score in sorted(scored, key=lambda item: item[1], reverse=True) if score > 0]
    if matched:
        return matched[:MAX_SOURCE_COUNT]
    if _terms(message) or _extract_requested_topic(message):
        return []
    return all_sources[:MAX_SOURCE_COUNT]


def _guess_knowledge_point(message: str, sources: list[KnowledgeSource]) -> str:
    requested_topic = _extract_requested_topic(message)
    if requested_topic:
        return requested_topic
    candidates = ["队列", "栈与队列", "循环队列", "链表", "二叉树", "机器学习", "Python"]
    for item in candidates:
        if item in message:
            return item
    for source in sources:
        points = _safe_json_list(source.knowledge_points)
        if points:
            return points[0]
    return "自主学习主题"


def _profile_focus(profile: dict[str, Any] | None) -> dict[str, list[str]]:
    if not isinstance(profile, dict):
        return {"weak_points": [], "frequent_errors": [], "recommendations": []}
    weak_points: list[str] = []
    for item in profile.get("knowledge_states", []):
        if not isinstance(item, dict):
            continue
        point = str(item.get("knowledge_point", "")).strip()
        state = str(item.get("state", "")).upper()
        mastery = item.get("mastery_score")
        is_weak = state in {"WEAK", "AT_RISK", "NEEDS_REVIEW"} or (isinstance(mastery, (int, float)) and mastery < 0.7)
        if point and is_weak:
            weak_points.append(point)

    frequent_errors = [
        str(item.get("label") or item.get("error_type") or "").strip()
        for item in profile.get("frequent_errors", [])
        if isinstance(item, dict) and str(item.get("label") or item.get("error_type") or "").strip()
    ]
    recommendations = [
        str(item.get("title") or item.get("suggested_action") or "").strip()
        for item in profile.get("recommendations", [])
        if isinstance(item, dict) and str(item.get("title") or item.get("suggested_action") or "").strip()
    ]
    return {
        "weak_points": weak_points[:3],
        "frequent_errors": frequent_errors[:3],
        "recommendations": recommendations[:2],
    }


def _fallback_slides(
    *,
    message: str,
    course: Course,
    knowledge_point: str,
    sources: list[KnowledgeSource],
    profile: dict[str, Any] | None = None,
) -> tuple[str, str, list[dict[str, Any]]]:
    title = f"{knowledge_point}讲解 PPT"
    source_ids = [source.id for source in sources[:3]]
    source_summary = sources[0].summary if sources else f"当前课程知识库未命中“{knowledge_point}”的直接引用，本页按用户请求生成通用学习框架。"
    focus = _profile_focus(profile)
    weak_text = "、".join(focus["weak_points"]) if focus["weak_points"] else "先按用户目标建立基础理解"
    error_text = "、".join(focus["frequent_errors"]) if focus["frequent_errors"] else "概念混淆、步骤遗漏、缺少自测"
    recommendation_text = "；".join(focus["recommendations"]) if focus["recommendations"] else f"完成后生成一组{knowledge_point}练习题做巩固"
    slides = [
        {
            "title": title,
            "subtitle": f"{course.name} · 自主学习资源",
            "bullets": [f"用户请求：{message}", f"学习目标：围绕{knowledge_point}形成可讲解、可练习、可复习的内容", "使用方式：预览后可加入资源中心并继续生成练习"],
            "speaker_notes": "开场先复述学生的实际请求，确认本资源服务于当前问题而不是固定模板。",
            "citation_ids": source_ids[:1],
            "layout": "cover",
        },
        {
            "title": f"{knowledge_point}的核心内容",
            "bullets": [source_summary, f"先回答“{knowledge_point}是什么、解决什么问题、怎么使用”", "把概念、示例和最小可运行片段放在同一条学习线上"],
            "speaker_notes": "如果知识库没有直接资料，讲解要明确这是按用户请求生成的通用学习内容，避免伪造课程引用。",
            "citation_ids": source_ids[:2],
            "layout": "content",
        },
        {
            "title": "按用户目标拆成学习路径",
            "bullets": [f"先梳理{knowledge_point}的基础规则和关键词", "再用 2 到 3 个小例子解释使用场景", "最后用练习题检查是否真的会迁移应用"],
            "speaker_notes": "这一页把学生的一句话需求拆成可执行学习步骤，避免只生成泛泛介绍。",
            "citation_ids": source_ids[:2],
            "layout": "content",
        },
        {
            "title": "结合学习画像的关注点",
            "bullets": [f"薄弱点参考：{weak_text}", f"常见错因参考：{error_text}", f"讲解{knowledge_point}时优先补足这些理解断点"],
            "speaker_notes": "画像只作为个性化适配依据，不把学生标签化；如果画像缺失，就按用户当前请求给出学习支架。",
            "citation_ids": source_ids[1:3] or source_ids[:1],
            "layout": "content",
        },
        {
            "title": "练习与自测设计",
            "bullets": [f"围绕{knowledge_point}设计基础识记题、代码阅读题和小任务", "每个练习都标注考查维度：概念、应用、易错点", "练习后把不确定点回填到 AI 对话继续追问"],
            "speaker_notes": "这一页用于把 PPT 从展示材料变成学习闭环的一部分。",
            "citation_ids": source_ids[:2],
            "layout": "content",
        },
        {
            "title": "学习总结",
            "bullets": [f"本 PPT 应始终围绕用户请求：{message}", f"下一步建议：{recommendation_text}", "保存资源后可继续生成笔记、卡片或练习题"],
            "speaker_notes": "收束重点，并引导学生把本次生成结果沉淀到资料库。",
            "citation_ids": source_ids[:1],
            "layout": "summary",
        },
    ]
    summary = f"围绕“{message}”生成 {len(slides)} 页 {knowledge_point}讲解 PPT。"
    return title, summary, slides


def _validate_model_slides(raw: dict[str, Any], allowed_source_ids: set[str]) -> tuple[str, str, list[dict[str, Any]]]:
    title = str(raw.get("title", "")).strip()[:120]
    summary = str(raw.get("summary", "")).strip()[:260]
    slides_raw = raw.get("slides")
    if not title or not isinstance(slides_raw, list) or len(slides_raw) < 4:
        raise ValueError("invalid ppt structure")
    slides: list[dict[str, Any]] = []
    for slide in slides_raw[:10]:
        if not isinstance(slide, dict):
            continue
        slide_title = str(slide.get("title", "")).strip()[:100]
        bullets = [str(item).strip()[:140] for item in slide.get("bullets", []) if str(item).strip()]
        notes = str(slide.get("speaker_notes", "")).strip()[:420]
        raw_citations = slide.get("citation_ids", [])
        citation_ids = [str(item) for item in raw_citations if str(item) in allowed_source_ids][:3] if isinstance(raw_citations, list) else []
        if slide_title and bullets:
            slides.append(
                {
                    "title": slide_title,
                    "subtitle": str(slide.get("subtitle", "")).strip()[:100],
                    "bullets": bullets[:5],
                    "speaker_notes": notes,
                    "citation_ids": citation_ids,
                    "layout": str(slide.get("layout", "content")).strip() or "content",
                }
            )
    if len(slides) < 4:
        raise ValueError("too few valid slides")
    return title, summary or f"{title}，共 {len(slides)} 页。", slides


async def _generate_model_slides(
    *,
    course: Course,
    message: str,
    knowledge_point: str,
    profile: dict[str, Any] | None,
    sources: list[KnowledgeSource],
) -> tuple[str, str, list[dict[str, Any]]] | None:
    settings = get_settings()
    if not settings.model_api_key or not settings.model_name:
        return None
    focus = _profile_focus(profile)
    payload = {
        "course": {"id": course.id, "name": course.name},
        "student_request": message,
        "request_contract": {
            "requested_topic": knowledge_point,
            "priority": [
                "必须优先满足 student_request 中的用户显式主题和目标。",
                "course 只提供当前学习环境，不能覆盖或替换用户显式主题。",
                "learner_profile 只用于个性化讲解顺序、薄弱点提醒和练习设计。",
                "knowledge_sources 只在与 requested_topic 直接相关时引用；不相关时 citation_ids 返回空数组。",
            ],
            "profile_focus": focus,
        },
        "knowledge_point": knowledge_point,
        "learner_profile": profile,
        "knowledge_sources": [_source_payload(source) for source in sources],
        "task": "生成一个可在浏览器预览并可导出为 PPTX 的完整中文教学演示文稿结构。",
        "requirements": [
            f"整份 PPT 必须围绕“{knowledge_point}”展开，标题、摘要或首页要点必须能看出该主题。",
            "首页要明确回应 student_request，而不是生成课程默认模板。",
            "结合 learner_profile 中的薄弱知识点、常见错因和推荐项安排讲解重点；没有画像时说明按当前请求组织。",
            "不是大纲，每页必须有标题、要点和讲稿提示。",
            "每页要点适合直接渲染成幻灯片。",
            "只能引用 knowledge_sources 中存在且与本页内容直接相关的 source_id。",
            "如果 knowledge_sources 为空或与主题不相关，所有 citation_ids 必须为空数组。",
            "不要生成虚构教材、论文或链接。",
        ],
        "output_schema": {
            "title": "PPT title",
            "summary": "short resource summary",
            "slides": "4-10 slides, each has title, subtitle, bullets, speaker_notes, citation_ids, layout",
        },
    }
    allowed = {source.id for source in sources}
    result = await chat_json(
        [
            {
                "role": "system",
                "content": "你是 CodeTrack 的教学资源生成器。只输出 JSON，不要 Markdown。",
            },
            {"role": "user", "content": _json_dumps(payload)},
        ],
        model=settings.model_name,
        api_key=settings.model_api_key,
        base_url=settings.model_api_base_url,
        timeout=35,
    )
    title, summary, slides = _validate_model_slides(result.data, allowed)
    _assert_slides_match_topic(title, summary, slides, knowledge_point)
    return title, summary, slides


def _add_text_box(slide, left, top, width, height, text: str, font_size, color, bold: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _render_pptx(resource_id: str, title: str, slides: list[dict[str, Any]], citations: list[dict[str, Any]]) -> str:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        raise ApiError(
            503,
            "PPT_RENDERER_NOT_INSTALLED",
            "PPT 生成依赖尚未安装，请先安装 backend/requirements.txt 中的 python-pptx。",
            details={"missing_module": exc.name},
        ) from exc

    settings = get_settings()
    storage_dir = Path(settings.resource_storage_dir) / "generated" / "ppt"
    storage_dir.mkdir(parents=True, exist_ok=True)
    path = storage_dir / f"{resource_id}.pptx"
    citation_titles = {item["source_id"]: item["title"] for item in citations}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, item in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 251, 255)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(23, 108, 245)
        bar.line.color.rgb = RGBColor(23, 108, 245)

        if index == 0:
            _add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.6), Inches(0.8), item["title"], Pt(34), RGBColor(15, 27, 69), True)
            subtitle = item.get("subtitle") or title
            _add_text_box(slide, Inches(0.86), Inches(2.25), Inches(9.8), Inches(0.45), subtitle, Pt(18), RGBColor(79, 91, 117))
            top = 3.15
        else:
            _add_text_box(slide, Inches(0.65), Inches(0.58), Inches(11.8), Inches(0.52), item["title"], Pt(25), RGBColor(15, 27, 69), True)
            top = 1.55

        bullet_box = slide.shapes.add_textbox(Inches(0.86), Inches(top), Inches(11.5), Inches(3.8))
        frame = bullet_box.text_frame
        frame.clear()
        for bullet_index, bullet in enumerate(item.get("bullets", [])):
            paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(18 if index else 16)
            paragraph.font.color.rgb = RGBColor(35, 48, 74)
            paragraph.space_after = Pt(10)

        notes = item.get("speaker_notes", "")
        if notes:
            _add_text_box(slide, Inches(0.86), Inches(5.6), Inches(8.4), Inches(0.42), f"讲稿提示：{notes}", Pt(10), RGBColor(100, 116, 139))

        ids = [source_id for source_id in item.get("citation_ids", []) if source_id in citation_titles]
        if ids:
            citation_text = "引用：" + "；".join(citation_titles[source_id] for source_id in ids[:2])
            _add_text_box(slide, Inches(0.86), Inches(6.55), Inches(10.6), Inches(0.28), citation_text, Pt(9), RGBColor(90, 105, 130))

        page = slide.shapes.add_textbox(Inches(11.9), Inches(6.55), Inches(0.8), Inches(0.28))
        page.text_frame.text = f"{index + 1}/{len(slides)}"
        page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        page.text_frame.paragraphs[0].font.size = Pt(9)
        page.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

    prs.save(path)
    return str(path)


def _resolve_libreoffice_command(settings: Any) -> str | None:
    configured = str(getattr(settings, "libreoffice_command", "") or "").strip()
    if configured:
        return configured
    for candidate in ("soffice", "libreoffice"):
        found = shutil.which(candidate)
        if found:
            return found
    return None


def _preview_base_metadata(renderer: str) -> dict[str, Any]:
    return {
        "preview_renderer": renderer,
        "preview_format": "PDF",
    }


def _preview_output_dir(resource_id: str) -> Path:
    return (Path(get_settings().resource_storage_dir) / "generated" / "previews" / resource_id).resolve()


def _power_shell_literal(value: str) -> str:
    return "'" + value.replace("'", "''") + "'"


def _render_ppt_pdf_preview_with_powerpoint(resource_id: str, source_path: Path, timeout_seconds: int) -> dict[str, Any]:
    metadata = _preview_base_metadata("powerpoint_pdf")
    if os.name != "nt":
        metadata["preview_error"] = "本机 PowerPoint 预览仅支持 Windows 环境。"
        return metadata

    power_shell = shutil.which("powershell") or shutil.which("pwsh")
    if not power_shell:
        metadata["preview_error"] = "未找到 PowerShell，无法调用本机 PowerPoint 生成预览。"
        return metadata

    preview_dir = _preview_output_dir(resource_id)
    preview_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = preview_dir / f"{source_path.stem}.pdf"
    if pdf_path.exists():
        pdf_path.unlink()

    script = f"""
$ErrorActionPreference = 'Stop'
$inputPath = {_power_shell_literal(str(source_path))}
$outputPath = {_power_shell_literal(str(pdf_path))}
$ppt = $null
$presentation = $null
try {{
  $ppt = New-Object -ComObject PowerPoint.Application
  $presentation = $ppt.Presentations.Open($inputPath, $true, $false, $false)
  $presentation.SaveAs($outputPath, 32)
}} finally {{
  if ($presentation -ne $null) {{
    $presentation.Close() | Out-Null
    [System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation) | Out-Null
  }}
  if ($ppt -ne $null) {{
    $ppt.Quit() | Out-Null
    [System.Runtime.InteropServices.Marshal]::ReleaseComObject($ppt) | Out-Null
  }}
  [GC]::Collect()
  [GC]::WaitForPendingFinalizers()
}}
""".strip()
    try:
        completed = subprocess.run(
            [power_shell, "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
            cwd=str(preview_dir),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_seconds,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        metadata["preview_error"] = f"PowerPoint 预览转换失败：{exc}"
        return metadata

    if completed.returncode != 0 or not pdf_path.exists() or pdf_path.stat().st_size < 100:
        detail = (completed.stderr or completed.stdout or "未生成 PDF 文件").strip()[-500:]
        metadata["preview_error"] = f"PowerPoint 预览转换失败：{detail}"
        return metadata

    metadata.update(
        {
            "preview_available": True,
            "preview_path": str(pdf_path),
            "preview_media_type": "application/pdf",
        }
    )
    return metadata


def _render_ppt_pdf_preview_with_libreoffice(resource_id: str, source_path: Path, timeout_seconds: int) -> dict[str, Any]:
    settings = get_settings()
    metadata = _preview_base_metadata("libreoffice_pdf")

    command = _resolve_libreoffice_command(settings)
    if not command:
        metadata["preview_error"] = "未找到 LibreOffice/soffice 命令，无法生成真实 PDF 预览。"
        return metadata

    preview_dir = _preview_output_dir(resource_id)
    preview_dir.mkdir(parents=True, exist_ok=True)
    expected_path = preview_dir / f"{source_path.stem}.pdf"
    if expected_path.exists():
        expected_path.unlink()

    args = [
        command,
        "--headless",
        "--invisible",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        str(preview_dir),
        str(source_path),
    ]
    try:
        completed = subprocess.run(
            args,
            cwd=str(preview_dir),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_seconds,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        metadata["preview_error"] = f"LibreOffice 预览转换失败：{exc}"
        return metadata

    pdf_path = expected_path if expected_path.exists() else None
    if not pdf_path:
        candidates = sorted(preview_dir.glob("*.pdf"), key=lambda item: item.stat().st_mtime, reverse=True)
        pdf_path = candidates[0] if candidates else None
    if completed.returncode != 0 or not pdf_path or not pdf_path.exists() or pdf_path.stat().st_size < 100:
        detail = (completed.stderr or completed.stdout or "未生成 PDF 文件").strip()[-500:]
        metadata["preview_error"] = f"LibreOffice 预览转换失败：{detail}"
        return metadata

    metadata.update(
        {
            "preview_available": True,
            "preview_path": str(pdf_path),
            "preview_media_type": "application/pdf",
        }
    )
    return metadata


def _render_ppt_pdf_preview(resource_id: str, pptx_path: str, file_format: str) -> dict[str, Any]:
    settings = get_settings()
    metadata = _preview_base_metadata("local_pdf")
    timeout_seconds = max(10, int(getattr(settings, "ppt_preview_timeout_seconds", 90) or 90))
    mode = str(getattr(settings, "ppt_preview_renderer", "auto") or "auto").strip().lower()
    aliases = {
        "": "auto",
        "office": "powerpoint",
        "powerpoint_pdf": "powerpoint",
        "libreoffice_pdf": "libreoffice",
        "soffice": "libreoffice",
    }
    mode = aliases.get(mode, mode)

    if not bool(getattr(settings, "ppt_preview_enabled", True)):
        metadata["preview_error"] = "PPT PDF 预览已关闭。"
        return metadata
    if file_format.upper() != "PPTX":
        metadata["preview_error"] = f"当前文件格式 {file_format} 不支持 PPT 预览转换。"
        return metadata

    source_path = Path(pptx_path).resolve()
    if not source_path.exists():
        metadata["preview_error"] = "PPTX 文件不存在，无法生成预览。"
        return metadata

    renderers = []
    if mode == "powerpoint":
        renderers = [_render_ppt_pdf_preview_with_powerpoint]
    elif mode == "libreoffice":
        renderers = [_render_ppt_pdf_preview_with_libreoffice]
    else:
        renderers = [_render_ppt_pdf_preview_with_powerpoint, _render_ppt_pdf_preview_with_libreoffice]

    errors: list[str] = []
    last_result: dict[str, Any] | None = None
    for renderer in renderers:
        result = renderer(resource_id, source_path, timeout_seconds)
        last_result = result
        if result.get("preview_available"):
            return result
        if result.get("preview_error"):
            errors.append(str(result["preview_error"]))

    if last_result:
        metadata.update({key: value for key, value in last_result.items() if key.startswith("preview_")})
    metadata["preview_error"] = "；".join(errors) or "未生成真实 PDF 预览。"
    return metadata


def _ppt_renderer_mode(settings: Any) -> str:
    mode = str(getattr(settings, "ppt_renderer", "auto") or "auto").strip().lower()
    aliases = {
        "": "auto",
        "python": "local",
        "python-pptx": "local",
        "local_pptx": "local",
        "pptmaster": "ppt_master",
        "ppt-master": "ppt_master",
    }
    return aliases.get(mode, mode)


def ppt_renderer_config_payload() -> dict[str, Any]:
    settings = get_settings()
    mode = _ppt_renderer_mode(settings)
    presenton_ready = presenton_configured(settings)
    ppt_master_bridge_ready = ppt_master_configured(settings)
    ppt_master_home = Path(settings.ppt_master_home).resolve() if settings.ppt_master_home else None
    ppt_master_skill_dir = ppt_master_home / "skills" / "ppt-master" if ppt_master_home else None
    ppt_master_official_ready = bool(
        ppt_master_bridge_ready
        and ppt_master_home
        and ppt_master_home.exists()
        and ppt_master_skill_dir
        and (ppt_master_skill_dir / "scripts" / "svg_to_pptx.py").exists()
    )
    if mode == "presenton":
        active = "presenton" if presenton_ready else "local_pptx"
    elif mode == "ppt_master":
        if ppt_master_official_ready:
            active = "ppt_master"
        elif ppt_master_bridge_ready:
            active = "ppt_master_bridge"
        else:
            active = "local_pptx"
    elif mode == "local":
        active = "local_pptx"
    elif mode == "auto":
        active = "presenton" if presenton_ready else "local_pptx"
    else:
        active = "local_pptx"
    return {
        "requested": mode,
        "active": active,
        "available": {
            "presenton": presenton_ready,
            "ppt_master": ppt_master_official_ready,
            "ppt_master_bridge": ppt_master_bridge_ready,
            "local_pptx": True,
        },
        "ppt_master": {
            "command_configured": bool(settings.ppt_master_command),
            "home_configured": bool(settings.ppt_master_home),
            "home_exists": bool(ppt_master_home and ppt_master_home.exists()),
            "skill_dir_exists": bool(ppt_master_skill_dir and ppt_master_skill_dir.exists()),
            "official_converter_exists": bool(
                ppt_master_skill_dir and (ppt_master_skill_dir / "scripts" / "svg_to_pptx.py").exists()
            ),
        },
        "fallback": active == "local_pptx" and mode not in {"local", "auto"},
    }


async def _render_ppt_resource_file(
    db: Session,
    state: PptResourceState,
    resource_id: str,
) -> tuple[str, str, dict[str, Any]]:
    settings = get_settings()
    mode = _ppt_renderer_mode(settings)
    metadata: dict[str, Any] = {"renderer": "local_pptx", "renderer_requested": mode}

    if mode not in {"auto", "presenton", "ppt_master", "local"}:
        metadata["renderer_config_error"] = f"未知 PPT 渲染器：{mode}"
        record_step(
            db,
            state["run"],
            step_name="select_ppt_renderer",
            step_order=5,
            status="FAILED",
            output_summary={"fallback": True, "reason": metadata["renderer_config_error"]},
        )

    should_try_ppt_master = mode == "ppt_master" and ppt_master_configured(settings)
    should_try_presenton = mode in {"auto", "presenton"} and presenton_configured(settings)

    if mode == "ppt_master" and not ppt_master_configured(settings):
        metadata["ppt_master_error"] = "PPT Master 未启用或未配置包装命令。"
        record_step(
            db,
            state["run"],
            step_name="select_ppt_master_renderer",
            step_order=5,
            status="FAILED",
            output_summary={"fallback": True, "reason": metadata["ppt_master_error"]},
        )

    if mode == "presenton" and not presenton_configured(settings):
        metadata["presenton_error"] = "Presenton 未启用或未配置服务地址。"
        record_step(
            db,
            state["run"],
            step_name="select_presenton_renderer",
            step_order=5,
            status="FAILED",
            output_summary={"fallback": True, "reason": metadata["presenton_error"]},
        )

    if should_try_ppt_master:
        output_dir = Path(settings.resource_storage_dir) / "generated" / "ppt_master"
        try:
            result = await generate_ppt_master_pptx(
                resource_id=resource_id,
                title=state["title"],
                message=state["message"],
                knowledge_point=state["knowledge_point"],
                slides=state["slides"],
                citations=state["citations"],
                output_dir=output_dir,
                settings=settings,
            )
        except PptMasterError as exc:
            metadata["ppt_master_error"] = str(exc)[:400]
            record_step(
                db,
                state["run"],
                step_name="render_ppt_master_pptx",
                step_order=5,
                status="FAILED",
                output_summary={"fallback": True, "reason": str(exc)[:240]},
            )
        else:
            provider_payload = result.get("provider_payload") if isinstance(result.get("provider_payload"), dict) else {}
            metadata.update(
                {
                    "renderer": "ppt_master",
                    "ppt_master_request_path": result.get("request_path"),
                    "ppt_master_stdout_tail": result.get("stdout_tail"),
                    "ppt_master_stderr_tail": result.get("stderr_tail"),
                    "ppt_master_project_id": provider_payload.get("project_id"),
                    "ppt_master_export_path": provider_payload.get("export_path"),
                    "ppt_master_implementation": provider_payload.get("implementation"),
                    "ppt_master_official_converter_error": provider_payload.get("official_converter_error"),
                }
            )
            record_step(
                db,
                state["run"],
                step_name="render_ppt_master_pptx",
                step_order=5,
                output_summary={"resource_id": resource_id, "file_format": result["file_format"], "renderer": "ppt_master"},
            )
            return str(result["file_path"]), str(result["file_format"]), metadata

    if should_try_presenton:
        output_dir = Path(settings.resource_storage_dir) / "generated" / "presenton"
        try:
            result = await generate_presenton_pptx(
                resource_id=resource_id,
                title=state["title"],
                message=state["message"],
                knowledge_point=state["knowledge_point"],
                slides=state["slides"],
                citations=state["citations"],
                output_dir=output_dir,
                settings=settings,
            )
        except PresentonError as exc:
            metadata["presenton_error"] = str(exc)[:400]
            record_step(
                db,
                state["run"],
                step_name="render_presenton_pptx",
                step_order=5,
                status="FAILED",
                output_summary={"fallback": True, "reason": str(exc)[:240]},
            )
        else:
            provider_payload = result.get("provider_payload") if isinstance(result.get("provider_payload"), dict) else {}
            metadata.update(
                {
                    "renderer": "presenton",
                    "renderer_requested": mode,
                    "presenton_presentation_id": provider_payload.get("presentation_id"),
                    "presenton_edit_path": provider_payload.get("edit_path"),
                    "presenton_edit_url": result.get("edit_url"),
                    "presenton_download_path": result.get("download_path"),
                    "presenton_download_url": result.get("download_url"),
                }
            )
            if isinstance(result.get("presenton_slides"), list):
                state["presenton_slides"] = result["presenton_slides"]
            record_step(
                db,
                state["run"],
                step_name="render_presenton_pptx",
                step_order=5,
                output_summary={
                    "resource_id": resource_id,
                    "file_format": result["file_format"],
                    "presentation_id": provider_payload.get("presentation_id"),
                },
            )
            return str(result["file_path"]), str(result["file_format"]), metadata

    path = _render_pptx(resource_id, state["title"], state["slides"], state["citations"])
    record_step(
        db,
        state["run"],
        step_name="render_local_pptx",
        step_order=6 if (metadata.get("presenton_error") or metadata.get("ppt_master_error") or metadata.get("renderer_config_error")) else 5,
        output_summary={"resource_id": resource_id, "file_format": "PPTX", "renderer": "local_pptx"},
    )
    return path, "PPTX", metadata


def _citation_titles(citations: list[dict[str, Any]]) -> str:
    titles = [str(item.get("title", "")).strip() for item in citations if str(item.get("title", "")).strip()]
    return "；".join(titles[:3]) or "课程知识库"


def _common_source_ids(sources: list[KnowledgeSource], limit: int = 3) -> list[str]:
    return [source.id for source in sources[:limit]]


def _mind_map_topic_family(knowledge_point: str) -> str:
    normalized = knowledge_point.lower()
    if re.search(r"java|python|语法|函数|方法|变量|列表|字典|字符串|编程", normalized):
        return "programming"
    if re.search(r"机器学习|监督学习|过拟合|正则化|梯度|损失函数|模型|训练集|验证集", normalized):
        return "machine_learning"
    return "data_structure"


def _display_knowledge_point(knowledge_point: str) -> str:
    return re.sub(r"java(?!script)", "Java", knowledge_point, flags=re.IGNORECASE)


def _fallback_document(
    *,
    message: str,
    course: Course,
    knowledge_point: str,
    sources: list[KnowledgeSource],
) -> tuple[str, str, dict[str, Any], int]:
    source_ids = _common_source_ids(sources)
    display_topic = _display_knowledge_point(knowledge_point)
    source_summary = sources[0].summary if sources else f"课程资料暂未直接命中“{display_topic}”，以下内容按用户问题和课程学习路径组织。"
    topic_family = _mind_map_topic_family(f"{message} {knowledge_point}")
    requested_need = "复习巩固"
    if re.search(r"入门|基础|小白|讲清楚|看不懂", message):
        requested_need = "从零理解"
    elif re.search(r"考试|测验|背|重点|复习", message):
        requested_need = "考前复习"
    elif re.search(r"代码|实现|编程|java|python|语法", message, flags=re.IGNORECASE):
        requested_need = "代码应用"
    elif re.search(r"对比|区别|关系", message):
        requested_need = "概念辨析"
    topic_domain = {
        "programming": "编程基础",
        "machine_learning": "机器学习",
        "data_structure": "数据结构",
    }.get(topic_family, course.name)

    if topic_family == "programming":
        rules_table = (
            "| 模块 | 必须掌握什么 | 判断自己会不会 |\n"
            "| --- | --- | --- |\n"
            f"| 变量与类型 | 知道{display_topic}里数据如何被声明、赋值和转换 | 能解释 int、double、boolean、String 等类型适合存什么 |\n"
            "| 表达式与运算符 | 理解算术、比较、逻辑运算的执行结果 | 能手算 `a > 3 && b != 0` 这类表达式 |\n"
            "| 分支结构 | 用 if / else 表达不同条件下的执行路径 | 能画出每个条件对应的分支 |\n"
            "| 循环结构 | 用 for / while 处理重复过程，并写清退出条件 | 能说出循环变量何时变化、何时停止 |\n"
            "| 方法调用 | 把重复逻辑封装成输入、处理、返回值 | 能区分参数、返回值和局部变量 |"
        )
        example_block = (
            "```java\n"
            "int score = 86;\n"
            "if (score >= 90) {\n"
            "    System.out.println(\"A\");\n"
            "} else if (score >= 60) {\n"
            "    System.out.println(\"Pass\");\n"
            "} else {\n"
            "    System.out.println(\"Retry\");\n"
            "}\n"
            "```\n"
            "读这段代码时不要先背语法名词，而是按三步走：先看变量保存了什么，再看条件从上到下如何判断，最后看实际输出落在哪个分支。"
        )
        mistake_table = (
            "| 常见问题 | 典型表现 | 修正方式 |\n"
            "| --- | --- | --- |\n"
            "| 只背语法格式 | 会写模板，但换个题就不知道变量怎么设 | 先用一句话写出输入、处理、输出 |\n"
            "| 条件顺序混乱 | 多个 if / else if 的结果和预期不一致 | 从最特殊条件写到最一般条件 |\n"
            "| 循环边界写错 | 少执行一次或多执行一次 | 写出第 1 次、最后 1 次循环变量的值 |\n"
            "| 类型混用 | 编译报错或结果被截断 | 每次赋值前确认左右两边类型是否匹配 |"
        )
        practice_list = (
            "- 用 3 个变量描述一个学生成绩判断场景。\n"
            "- 写一段 if / else，把分数映射成 A、B、C、D。\n"
            "- 写一个 for 循环输出 1 到 10 中的偶数。\n"
            "- 解释你的循环为什么不会多输出或少输出。"
        )
    elif topic_family == "machine_learning":
        rules_table = (
            "| 模块 | 必须掌握什么 | 判断自己会不会 |\n"
            "| --- | --- | --- |\n"
            f"| 问题定义 | 知道{display_topic}要预测、分类或解释什么 | 能说清输入特征和目标变量 |\n"
            "| 数据划分 | 区分训练集、验证集和测试集 | 能解释为什么不能只看训练集效果 |\n"
            "| 模型训练 | 理解参数如何根据损失函数被调整 | 能说出损失下降代表什么 |\n"
            "| 泛化评估 | 判断模型是否过拟合或欠拟合 | 能根据训练/验证表现提出改进方向 |"
        )
        example_block = (
            "| 观察现象 | 可能问题 | 下一步 |\n"
            "| --- | --- | --- |\n"
            "| 训练准确率很高，验证准确率低 | 过拟合 | 减少复杂度、增加正则化或补充数据 |\n"
            "| 训练和验证准确率都低 | 欠拟合 | 增加特征、换模型或训练更充分 |\n"
            "| 指标波动大 | 数据量或划分不稳定 | 做交叉验证或检查样本分布 |"
        )
        mistake_table = (
            "| 常见问题 | 典型表现 | 修正方式 |\n"
            "| --- | --- | --- |\n"
            "| 把测试集当验证集反复调参 | 最终指标虚高 | 测试集只在最后使用 |\n"
            "| 只看准确率 | 类别不均衡时误判模型好坏 | 同时看召回率、精确率、F1 |\n"
            "| 忽略数据泄漏 | 线上效果明显变差 | 检查特征是否包含未来信息 |"
        )
        practice_list = (
            "- 写出一个机器学习任务的输入、输出和评价指标。\n"
            "- 说明训练集、验证集、测试集各自用途。\n"
            "- 给出一个过拟合现象，并写出两种处理办法。"
        )
    else:
        rules_table = (
            "| 模块 | 必须掌握什么 | 判断自己会不会 |\n"
            "| --- | --- | --- |\n"
            f"| 抽象结构 | 知道{display_topic}存储什么、限制什么 | 能画出一次操作前后的状态 |\n"
            "| 基本操作 | 掌握插入、删除、查找或访问规则 | 能手推连续 3 次操作后的结果 |\n"
            "| 边界条件 | 处理空结构、满结构、首尾位置 | 能写出最小用例和边界用例 |\n"
            "| 复杂度 | 知道关键操作的时间/空间代价 | 能解释为什么某个操作快或慢 |"
        )
        example_block = (
            "| 操作 | 操作前 | 操作后 | 需要检查 |\n"
            "| --- | --- | --- | --- |\n"
            f"| 插入元素 | {display_topic}当前状态 | 元素进入指定位置 | 是否为空、是否满、指针/下标是否更新 |\n"
            "| 删除元素 | 至少存在一个元素 | 目标元素被移除 | 删除后结构是否仍然连通或有序 |\n"
            "| 访问元素 | 给定位置或条件 | 返回对应元素 | 越界、空结构、重复元素如何处理 |"
        )
        mistake_table = (
            "| 常见问题 | 典型表现 | 修正方式 |\n"
            "| --- | --- | --- |\n"
            "| 只记住普通情况 | 一遇到空结构就错 | 每个操作先问能不能做，再问怎么做 |\n"
            "| 更新顺序错误 | 指针、下标或计数不一致 | 画出操作前后两张状态图 |\n"
            "| 忘记复杂度 | 会写代码但不会解释效率 | 把循环次数和数据规模对应起来 |"
        )
        practice_list = (
            f"- 画出{display_topic}执行 3 次连续操作后的状态变化。\n"
            "- 写出一个普通用例、一个最小用例、一个边界用例。\n"
            "- 解释每个核心操作的时间复杂度。"
        )

    title = f"{display_topic}学习讲义"
    sections = [
        {
            "heading": "用户需求拆解",
            "paragraphs": [
                f"用户输入是“{message}”，核心不是简单要一份资料，而是希望把“{display_topic}”细化成能理解、能复习、能练习的学习材料。",
                f"- 主题域：{topic_domain}\n- 关联课程：{course.name}\n- 需求类型：{requested_need}\n- 建议阅读方式：先看框架，再看例子，最后做自检任务。",
            ],
            "citation_ids": source_ids[:1],
        },
        {
            "heading": "一句话框架",
            "paragraphs": [
                source_summary,
                f"学习“{display_topic}”时，不要只记定义，要同时抓住三件事：它解决什么问题、规则如何运行、遇到边界情况时怎么判断。"
            ],
            "citation_ids": source_ids[:2],
        },
        {
            "heading": "必须掌握的规则",
            "paragraphs": [
                rules_table,
                "这张表可以作为预习和复习清单：如果某一行不能用自己的话解释，就说明这一块还没有真正掌握。"
            ],
            "citation_ids": source_ids[:2],
        },
        {
            "heading": "带着例子走一遍",
            "paragraphs": [
                example_block,
                "看例子时建议把每一步的输入、状态变化和输出写出来，这比只看最终答案更能发现薄弱点。"
            ],
            "citation_ids": source_ids[:2],
        },
        {
            "heading": "高频易错点",
            "paragraphs": [
                mistake_table,
                "如果你发现自己错在同一类问题上，可以把这一行单独保存成错因标签，后续让 AI 导师按这个错因继续出题。"
            ],
            "citation_ids": source_ids[1:3] or source_ids[:1],
        },
        {
            "heading": "10 分钟自检任务",
            "paragraphs": [
                practice_list,
                "自检时不要只写答案，还要写一句“为什么这样做”。能解释原因，才算真正掌握。"
            ],
            "citation_ids": source_ids[:1],
        },
        {
            "heading": "下一步学习建议",
            "paragraphs": [
                f"如果本文档读完仍觉得抽象，下一步建议继续生成“{display_topic}练习题”或“{display_topic}思维导图”，用题目和结构图反向检验理解。",
                f"推荐追问：请围绕{display_topic}出 5 道由浅入深的题，并标注每题考查的知识点和易错点。"
            ],
            "citation_ids": source_ids[:1],
        },
    ]
    metadata = {
        "document_style": "study-handout",
        "document_depth": "expanded",
        "requested_need": requested_need,
        "topic_family": topic_family,
        "topic_domain": topic_domain,
    }
    summary = f"已将“{message}”细化为 {len(sections)} 节学习讲义，包含规则表、示例、易错点和自检任务。"
    return title, summary, {"sections": sections, "metadata": metadata}, len(sections)


def _fallback_mind_map(
    *,
    message: str,
    course: Course,
    knowledge_point: str,
    sources: list[KnowledgeSource],
    profile: dict[str, Any] | None = None,
) -> tuple[str, str, dict[str, Any], int]:
    source_ids = _common_source_ids(sources)
    focus = _profile_focus(profile)
    topic_family = _mind_map_topic_family(knowledge_point)
    profile_matches_topic = any(
        knowledge_point in item or item in knowledge_point
        for item in [*focus["weak_points"], *focus["frequent_errors"]]
        if item
    )
    first_source_summary = sources[0].summary if sources else f"当前课程资料未覆盖“{knowledge_point}”，本导图按用户主题生成基础学习框架。"
    if topic_family == "programming":
        weak_text = "、".join(focus["weak_points"]) if profile_matches_topic and focus["weak_points"] else "变量类型、流程控制、函数调用"
        error_text = "、".join(focus["frequent_errors"]) if profile_matches_topic and focus["frequent_errors"] else "类型不匹配、作用域混淆、语法符号遗漏"
        branch_labels = [
            {
                "id": "concept",
                "title": "语言基础",
                "node_type": "concept",
                "relationship_type": "contains",
                "children": [
                    ("变量与类型", "理解变量声明、基本类型、引用类型和类型转换。", "concept", "contains"),
                    ("表达式", "掌握运算符优先级、布尔表达式和赋值表达式。", "concept", "contains"),
                    ("代码规范", "用命名、缩进和注释提升可读性。", "practice", "solves"),
                ],
            },
            {
                "id": "procedure",
                "title": "流程控制",
                "node_type": "procedure",
                "relationship_type": "contains",
                "children": [
                    ("条件分支", "用 if / else 表达不同条件下的执行路径。", "procedure", "contains"),
                    ("循环结构", "用 for / while 处理重复任务，并明确退出条件。", "procedure", "causes"),
                    ("输入输出", "把读入、处理和输出拆成清晰步骤。", "procedure", "next_step"),
                ],
            },
            {
                "id": "structure",
                "title": "函数与对象",
                "node_type": "concept",
                "relationship_type": "contains",
                "children": [
                    ("函数/方法", "用参数、返回值和局部变量封装一段逻辑。", "concept", "contains"),
                    ("数组或集合", "选择合适的数据容器保存多项数据。", "concept", "example_of"),
                    ("类与对象", "理解属性、方法和对象实例之间的关系。", "concept", "example_of"),
                ],
            },
            {
                "id": "mistake",
                "title": "常见错误",
                "node_type": "mistake",
                "relationship_type": "common_mistake",
                "children": [
                    ("类型不匹配", "检查赋值、比较和方法参数类型是否一致。", "mistake", "common_mistake"),
                    ("作用域混淆", "确认变量在哪个代码块中声明和使用。", "mistake", "causes"),
                    ("符号遗漏", "重点检查分号、括号、引号和缩进层级。", "mistake", "common_mistake"),
                ],
            },
            {
                "id": "profile",
                "title": "画像提醒",
                "node_type": "profile_tip",
                "relationship_type": "common_mistake",
                "children": [
                    ("薄弱点", weak_text, "profile_tip", "common_mistake"),
                    ("高频错因", error_text, "profile_tip", "causes"),
                    ("复习策略", "先读懂小程序，再手写一个同类型小例子。", "next_action", "solves"),
                ],
            },
            {
                "id": "practice",
                "title": "练习路径",
                "node_type": "practice",
                "relationship_type": "next_step",
                "children": [
                    ("读代码", "先预测输出，再运行验证自己的理解。", "practice", "next_step"),
                    ("改错题", "从编译错误和运行结果反推语法问题。", "practice", "solves"),
                    ("小程序", f"围绕{knowledge_point}写一个 20 行以内的练习。", "next_action", "next_step"),
                ],
            },
        ]
    elif topic_family == "machine_learning":
        weak_text = "、".join(focus["weak_points"]) if focus["weak_points"] else f"{knowledge_point}的指标理解和方法选择"
        error_text = "、".join(focus["frequent_errors"]) if focus["frequent_errors"] else "概念混淆、指标误读、训练验证流程不清"
        branch_labels = [
            {
                "id": "concept",
                "title": "核心概念",
                "node_type": "concept",
                "relationship_type": "contains",
                "children": [
                    (f"{knowledge_point}是什么", "先明确概念定义、问题背景和适用前提。", "concept", "contains"),
                    ("课程依据", first_source_summary, "concept", "example_of"),
                    ("关键术语", "把数据、模型、损失、指标和泛化放到同一框架理解。", "concept", "contains"),
                ],
            },
            {
                "id": "procedure",
                "title": "建模流程",
                "node_type": "procedure",
                "relationship_type": "contains",
                "children": [
                    ("数据划分", "区分训练集、验证集和测试集的作用。", "procedure", "prerequisite"),
                    ("训练评估", "用指标观察模型是否学到可泛化规律。", "procedure", "causes"),
                    ("调参与改进", "根据误差表现选择正则化、特征或模型调整。", "procedure", "solves"),
                ],
            },
            {
                "id": "mistake",
                "title": "常见误区",
                "node_type": "mistake",
                "relationship_type": "common_mistake",
                "children": [
                    ("只看训练效果", "训练集表现好不代表泛化能力好。", "mistake", "common_mistake"),
                    ("指标误读", "不同任务要选择匹配的评价指标。", "mistake", "causes"),
                    ("泄漏数据", "训练阶段不能提前看到测试信息。", "mistake", "common_mistake"),
                ],
            },
            {
                "id": "profile",
                "title": "画像提醒",
                "node_type": "profile_tip",
                "relationship_type": "common_mistake",
                "children": [
                    ("薄弱点", weak_text, "profile_tip", "common_mistake"),
                    ("高频错因", error_text, "profile_tip", "causes"),
                    ("复习策略", f"围绕{knowledge_point}先画流程，再做概念判断题。", "next_action", "solves"),
                ],
            },
            {
                "id": "practice",
                "title": "练习路径",
                "node_type": "practice",
                "relationship_type": "next_step",
                "children": [
                    ("概念辨析", "用判断题检查术语边界。", "practice", "next_step"),
                    ("案例分析", "根据训练/验证表现判断问题原因。", "practice", "example_of"),
                    ("生成练习", f"保存导图后生成一组{knowledge_point}练习题。", "next_action", "next_step"),
                ],
            },
        ]
    else:
        weak_text = "、".join(focus["weak_points"]) if focus["weak_points"] else f"{knowledge_point}的边界和迁移应用"
        error_text = "、".join(focus["frequent_errors"]) if focus["frequent_errors"] else "概念混淆、更新顺序、边界遗漏"
        branch_labels = [
            {
                "id": "concept",
                "title": "核心概念",
                "node_type": "concept",
                "relationship_type": "contains",
                "children": [
                    (f"{knowledge_point}是什么", "先说清定义、解决的问题和使用场景。", "concept", "contains"),
                    ("课程依据", first_source_summary, "concept", "example_of"),
                    ("适用场景", f"判断{knowledge_point}适合解决哪类输入、状态或结构问题。", "example", "example_of"),
                ],
            },
            {
                "id": "procedure",
                "title": "关键过程",
                "node_type": "procedure",
                "relationship_type": "contains",
                "children": [
                    ("输入与输出", "明确入口条件、输出目标和中间状态。", "procedure", "contains"),
                    ("状态变化", "跟踪每一步操作后结构或变量如何变化。", "procedure", "causes"),
                    ("复杂度关注", "理解时间复杂度和空间开销来自哪些步骤。", "concept", "contains"),
                ],
            },
            {
                "id": "boundary",
                "title": "边界条件",
                "node_type": "mistake",
                "relationship_type": "common_mistake",
                "children": [
                    ("空结构", "先判断空输入，避免访问不存在的节点或元素。", "mistake", "common_mistake"),
                    ("最小用例", "只含一个元素时，返回值和状态更新最容易暴露问题。", "mistake", "common_mistake"),
                    ("连续操作", "多次插入、删除或遍历后检查状态是否仍一致。", "practice", "solves"),
                ],
            },
            {
                "id": "profile",
                "title": "画像提醒",
                "node_type": "profile_tip",
                "relationship_type": "common_mistake",
                "children": [
                    ("薄弱点", weak_text, "profile_tip", "common_mistake"),
                    ("高频错因", error_text, "profile_tip", "causes"),
                    ("复习策略", f"围绕{knowledge_point}先画过程，再做边界自测。", "next_action", "solves"),
                ],
            },
            {
                "id": "practice",
                "title": "练习路径",
                "node_type": "practice",
                "relationship_type": "next_step",
                "children": [
                    ("画状态图", "用小样例手动画出每一步变化。", "practice", "next_step"),
                    ("写伪过程", "先写步骤和判断条件，再补具体代码。", "practice", "next_step"),
                    ("生成练习", f"保存导图后生成一组{knowledge_point}练习题。", "next_action", "next_step"),
                ],
            },
        ]
    center_id = "center"
    center_summary = f"{course.name} · {message}" if source_ids else f"当前课程资料未覆盖该主题 · {message}"
    nodes = [
        {
            "id": center_id,
            "node_id": center_id,
            "parent_id": None,
            "label": f"{knowledge_point}学习地图",
            "title": f"{knowledge_point}学习地图",
            "level": 0,
            "depth": 0,
            "summary": center_summary,
            "node_type": "central_topic",
            "knowledge_points": [knowledge_point],
            "citation_ids": source_ids[:1],
            "citations": source_ids[:1],
            "confidence": 0.88 if source_ids else 0.58,
        }
    ]
    edges: list[dict[str, Any]] = []
    for index, branch in enumerate(branch_labels, start=1):
        node_id = f"branch_{branch['id']}"
        branch_source_ids = source_ids[: min(len(source_ids), 2)]
        nodes.append(
            {
                "id": node_id,
                "node_id": node_id,
                "parent_id": center_id,
                "label": branch["title"],
                "title": branch["title"],
                "level": 1,
                "depth": 1,
                "summary": " / ".join(str(child[0]) for child in branch["children"]),
                "node_type": branch["node_type"],
                "knowledge_points": [knowledge_point],
                "citation_ids": branch_source_ids,
                "citations": branch_source_ids,
                "confidence": 0.86 if branch_source_ids else 0.56,
            }
        )
        edges.append(
            {
                "source": center_id,
                "target": node_id,
                "source_node_id": center_id,
                "target_node_id": node_id,
                "relationship_type": branch["relationship_type"],
                "label": "展开",
            }
        )
        for child_index, (child_title, child_summary, child_type, child_relation) in enumerate(branch["children"], start=1):
            child_id = f"{branch['id']}_{child_index}"
            child_source_ids = source_ids[index % len(source_ids): index % len(source_ids) + 1] if source_ids else []
            nodes.append(
                {
                    "id": child_id,
                    "node_id": child_id,
                    "parent_id": node_id,
                    "label": child_title,
                    "title": child_title,
                    "level": 2,
                    "depth": 2,
                    "summary": child_summary,
                    "node_type": child_type,
                    "knowledge_points": [knowledge_point],
                    "citation_ids": child_source_ids,
                    "citations": child_source_ids,
                    "confidence": 0.84 if child_source_ids else 0.54,
                }
            )
            edges.append(
                {
                    "source": node_id,
                    "target": child_id,
                    "source_node_id": node_id,
                    "target_node_id": child_id,
                    "relationship_type": child_relation,
                    "label": "包含",
                }
            )
    cross_edges = [
        ("branch_profile", "practice_1", "solves", "转为练习"),
    ]
    if topic_family == "programming":
        cross_edges.extend(
            [
                ("procedure_1", "mistake_2", "causes", "分支变量影响作用域"),
                ("mistake_1", "practice_2", "solves", "用改错题巩固"),
            ]
        )
    elif topic_family == "machine_learning":
        cross_edges.extend(
            [
                ("procedure_2", "mistake_1", "causes", "评估表现暴露问题"),
                ("mistake_2", "practice_1", "solves", "用辨析题巩固"),
            ]
        )
    else:
        cross_edges.extend(
            [
                ("branch_boundary", "profile_2", "causes", "关联错因"),
                ("procedure_2", "boundary_3", "causes", "连续状态影响边界"),
                ("boundary_2", "practice_2", "solves", "用伪过程检查"),
            ]
        )
    edges.extend(
        [
            {
                "source": source_id,
                "target": target_id,
                "source_node_id": source_id,
                "target_node_id": target_id,
                "relationship_type": relation,
                "label": label,
            }
            for source_id, target_id, relation, label in cross_edges
        ]
    )
    title = f"{knowledge_point}思维导图"
    summary = f"围绕“{message}”生成 {len(nodes)} 个节点的学习地图。"
    risk_flags = [] if source_ids else ["当前课程资料未覆盖该主题", "引用不足"]
    return (
        title,
        summary,
        {
            "artifact_type": "MIND_MAP",
            "central_topic": f"{knowledge_point}学习地图",
            "nodes": nodes,
            "edges": edges,
            "risk_flags": risk_flags,
            "recommended_next_actions": ["生成练习", "保存到资源中心", "继续追问"],
            "metadata": {
                "planner": "rule_fallback_mind_map_planner",
                "topic_family": topic_family,
                "relationship_types": [
                    "prerequisite",
                    "contains",
                    "causes",
                    "solves",
                    "example_of",
                    "common_mistake",
                    "next_step",
                    "contrast",
                ],
            },
        },
        len(nodes),
    )


def _fallback_practice_set(
    *,
    message: str,
    knowledge_point: str,
    sources: list[KnowledgeSource],
) -> tuple[str, str, dict[str, Any], int]:
    source_ids = _common_source_ids(sources)
    title = f"{knowledge_point}巩固练习"
    questions = [
        {
            "type": "single_choice",
            "stem": f"关于{knowledge_point}的核心规则，下列说法更合理的是哪一项？",
            "options": ["只关注最终结果", "需要跟踪状态变化和边界情况", "不需要测试空输入", "实现时无需考虑复杂度"],
            "answer": "需要跟踪状态变化和边界情况",
            "analysis": f"{knowledge_point}学习不能只记结论，必须能解释状态如何变化。",
            "citation_ids": source_ids[:1],
        },
        {
            "type": "short_answer",
            "stem": f"请用自己的话说明{knowledge_point}最容易出错的一个边界情况。",
            "answer": "示例：空结构、满结构、只有一个元素时的状态更新。",
            "analysis": "能主动识别边界情况，说明已经从概念理解走向实现检查。",
            "citation_ids": source_ids[:2],
        },
        {
            "type": "process",
            "stem": f"给定一个小规模样例，请手动画出{knowledge_point}每一步操作后的状态。",
            "answer": "按操作顺序列出状态变量或结构内容，确保每一步都符合不变式。",
            "analysis": "过程追踪可以发现只看最终结果时漏掉的更新顺序问题。",
            "citation_ids": source_ids[:2],
        },
        {
            "type": "debug",
            "stem": f"如果{knowledge_point}实现只能通过普通用例，却过不了边界用例，你会优先检查哪里？",
            "answer": "优先检查空/满判断、指针或索引更新顺序、返回值约定。",
            "analysis": "边界失败通常不是整体思路错误，而是状态维护细节不完整。",
            "citation_ids": source_ids[1:3] or source_ids[:1],
        },
        {
            "type": "reflection",
            "stem": f"完成“{message}”后，写出一个你还不确定的点。",
            "answer": "把不确定点带回 AI 对话继续追问，或生成配套讲解资源。",
            "analysis": "反思题用于把一次练习转成下一轮学习目标。",
            "citation_ids": source_ids[:1],
        },
    ]
    summary = f"围绕“{message}”生成 {len(questions)} 道巩固练习。"
    return title, summary, {"questions": questions}, len(questions)


def _fallback_knowledge_card(
    *,
    message: str,
    knowledge_point: str,
    sources: list[KnowledgeSource],
) -> tuple[str, str, dict[str, Any], int]:
    source_ids = _common_source_ids(sources)
    cards = [
        {
            "front": f"{knowledge_point}是什么？",
            "back": f"{knowledge_point}是一类需要理解规则、状态和边界条件的学习主题。",
            "tips": ["先讲规则", "再画状态", "最后测边界"],
            "citation_ids": source_ids[:1],
        },
        {
            "front": "为什么要画过程？",
            "back": "过程图能暴露状态变化是否符合不变式，特别适合排查更新顺序问题。",
            "tips": ["记录每一步", "标出入口出口", "关注最小用例"],
            "citation_ids": source_ids[:2],
        },
        {
            "front": "如何自检？",
            "back": "用普通用例、空/满或最小用例、连续操作用例分别验证。",
            "tips": ["不要只测一个例子", "把失败原因写成一句话"],
            "citation_ids": source_ids[1:3] or source_ids[:1],
        },
    ]
    title = f"{knowledge_point}知识卡片"
    summary = f"围绕“{message}”生成 {len(cards)} 张复习卡片。"
    return title, summary, {"cards": cards}, len(cards)


def _fallback_podcast_script(
    *,
    message: str,
    course: Course,
    knowledge_point: str,
    sources: list[KnowledgeSource],
) -> tuple[str, str, dict[str, Any], int]:
    source_ids = _common_source_ids(sources)
    title = f"{knowledge_point}播客讲解稿"
    segments = [
        {
            "speaker": "主持人",
            "label": "开场",
            "text": f"今天我们用一段短讲，梳理{course.name}里的{knowledge_point}。问题来自：{message}。",
            "citation_ids": source_ids[:1],
        },
        {
            "speaker": "讲解者",
            "label": "概念",
            "text": f"先抓住主线：{knowledge_point}不是孤立术语，要和操作规则、状态变化、边界情况一起理解。",
            "citation_ids": source_ids[:2],
        },
        {
            "speaker": "主持人",
            "label": "追问",
            "text": "那学生最容易卡在哪里？",
            "citation_ids": [],
        },
        {
            "speaker": "讲解者",
            "label": "易错点",
            "text": "常见卡点是只记住普通情况，没有检查空结构、满结构、最小规模和连续操作后的状态。",
            "citation_ids": source_ids[1:3] or source_ids[:1],
        },
        {
            "speaker": "主持人",
            "label": "行动",
            "text": "听完之后，建议先画一个小例子的状态变化，再生成配套练习检验自己。",
            "citation_ids": source_ids[:1],
        },
    ]
    summary = f"围绕“{message}”生成 {len(segments)} 段双人播客讲解稿。"
    return title, summary, {"segments": segments}, len(segments)


def _fallback_resource_payload(
    *,
    resource_type: str,
    message: str,
    course: Course,
    knowledge_point: str,
    sources: list[KnowledgeSource],
) -> tuple[str, str, dict[str, Any], int]:
    if resource_type == "DOCUMENT":
        return _fallback_document(message=message, course=course, knowledge_point=knowledge_point, sources=sources)
    if resource_type == "MIND_MAP":
        return _fallback_mind_map(message=message, course=course, knowledge_point=knowledge_point, sources=sources)
    if resource_type == "PRACTICE_SET":
        return _fallback_practice_set(message=message, knowledge_point=knowledge_point, sources=sources)
    if resource_type == "KNOWLEDGE_CARD":
        return _fallback_knowledge_card(message=message, knowledge_point=knowledge_point, sources=sources)
    if resource_type == "PODCAST_SCRIPT":
        return _fallback_podcast_script(message=message, course=course, knowledge_point=knowledge_point, sources=sources)
    raise ApiError(400, "UNSUPPORTED_RESOURCE_TYPE", "暂不支持该资源类型。", details={"resource_type": resource_type})


def _render_docx(resource_id: str, title: str, payload: dict[str, Any], citations: list[dict[str, Any]]) -> tuple[str, str]:
    try:
        from docx import Document
    except ModuleNotFoundError as exc:
        raise ApiError(
            503,
            "DOC_RENDERER_NOT_INSTALLED",
            "文档生成依赖尚未安装，请先安装 backend/requirements.txt 中的 python-docx。",
            details={"missing_module": exc.name},
        ) from exc

    settings = get_settings()
    storage_dir = Path(settings.resource_storage_dir) / "generated" / "document"
    storage_dir.mkdir(parents=True, exist_ok=True)
    path = storage_dir / f"{resource_id}.docx"
    document = Document()
    document.add_heading(title, level=1)

    def add_markdownish_table(lines: list[str]) -> bool:
        if len(lines) < 2:
            return False
        divider = lines[1].replace("|", "").replace("-", "").replace(":", "").strip()
        if divider:
            return False
        rows = [[cell.strip() for cell in line.strip().strip("|").split("|")] for line in lines if line.strip()]
        rows = [row for index, row in enumerate(rows) if index != 1 and any(row)]
        if not rows:
            return False
        table = document.add_table(rows=len(rows), cols=max(len(row) for row in rows))
        try:
            table.style = "Table Grid"
        except KeyError:
            pass
        for row_index, row in enumerate(rows):
            for col_index, cell in enumerate(row):
                table.rows[row_index].cells[col_index].text = cell
        return True

    def add_markdownish_paragraphs(paragraphs: list[Any]) -> None:
        for raw in paragraphs:
            text = str(raw).strip()
            if not text:
                continue
            lines = [line.rstrip() for line in text.splitlines() if line.strip()]
            if lines and lines[0].startswith("|") and add_markdownish_table(lines):
                continue
            in_code = False
            for line in lines:
                stripped = line.strip()
                if stripped.startswith("```"):
                    in_code = not in_code
                    continue
                if stripped.startswith("- "):
                    document.add_paragraph(stripped[2:].strip(), style="List Bullet")
                elif stripped.startswith("### "):
                    document.add_heading(stripped[4:].strip(), level=3)
                elif in_code:
                    paragraph = document.add_paragraph()
                    run = paragraph.add_run(stripped)
                    run.font.name = "Consolas"
                else:
                    document.add_paragraph(stripped)

    for section in payload.get("sections", []):
        if not isinstance(section, dict):
            continue
        document.add_heading(str(section.get("heading", "学习小节")), level=2)
        add_markdownish_paragraphs(section.get("paragraphs", []))
    if citations:
        document.add_heading("引用来源", level=2)
        for citation in citations[:5]:
            document.add_paragraph(f"{citation.get('title', '')}：{citation.get('summary', '')}", style=None)
    document.save(path)
    return str(path), "DOCX"


def _markdown_lines(resource_type: str, title: str, payload: dict[str, Any], citations: list[dict[str, Any]]) -> list[str]:
    lines = [f"# {title}", ""]
    if resource_type == "DOCUMENT":
        for section in payload.get("sections", []):
            if isinstance(section, dict):
                lines.extend([f"## {section.get('heading', '学习小节')}", ""])
                for paragraph in section.get("paragraphs", []):
                    if str(paragraph).strip():
                        lines.extend([str(paragraph).strip(), ""])
    elif resource_type == "MIND_MAP":
        nodes = payload.get("nodes", [])
        edges = payload.get("edges", [])
        lines.extend(["## 节点", ""])
        for node in nodes if isinstance(nodes, list) else []:
            if isinstance(node, dict):
                indent = "  " * int(node.get("level", 0) or 0)
                confidence = node.get("confidence")
                confidence_text = f"；置信度 {round(float(confidence) * 100)}%" if isinstance(confidence, (int, float)) else ""
                node_type = f"；类型 {node.get('node_type')}" if node.get("node_type") else ""
                citation_ids = node.get("citation_ids") if isinstance(node.get("citation_ids"), list) else []
                citation_text = f"；引用 {', '.join(str(item) for item in citation_ids)}" if citation_ids else ""
                lines.append(f"{indent}- {node.get('label') or node.get('title', '')}：{node.get('summary', '')}{node_type}{confidence_text}{citation_text}")
        lines.extend(["", "## 关系", ""])
        for edge in edges if isinstance(edges, list) else []:
            if isinstance(edge, dict):
                relation = edge.get("relationship_type", "contains")
                lines.append(f"- {edge.get('source')} -> {edge.get('target')}：{edge.get('label', '')}（{relation}）")
    elif resource_type == "PRACTICE_SET":
        for index, question in enumerate(payload.get("questions", []), start=1):
            if not isinstance(question, dict):
                continue
            lines.extend([f"## 第 {index} 题", "", str(question.get("stem", "")), ""])
            options = question.get("options", [])
            if isinstance(options, list):
                lines.extend([f"- {option}" for option in options])
                lines.append("")
            lines.extend([f"答案：{question.get('answer', '')}", "", f"解析：{question.get('analysis', '')}", ""])
    elif resource_type == "KNOWLEDGE_CARD":
        for index, card in enumerate(payload.get("cards", []), start=1):
            if isinstance(card, dict):
                lines.extend([f"## 卡片 {index}", "", f"正面：{card.get('front', '')}", "", f"背面：{card.get('back', '')}", ""])
                tips = card.get("tips", [])
                if isinstance(tips, list):
                    lines.extend([f"- {tip}" for tip in tips])
                    lines.append("")
    elif resource_type == "PODCAST_SCRIPT":
        for segment in payload.get("segments", []):
            if isinstance(segment, dict):
                lines.extend([f"## {segment.get('speaker', '')} · {segment.get('label', '')}", "", str(segment.get("text", "")), ""])
    else:
        lines.extend([json.dumps(payload, ensure_ascii=False, indent=2), ""])
    if citations:
        lines.extend(["## 引用来源", ""])
        for citation in citations[:5]:
            lines.append(f"- {citation.get('title', '')}：{citation.get('summary', '')}")
    return lines


def _attach_document_markdown_preview(title: str, payload: dict[str, Any], citations: list[dict[str, Any]]) -> dict[str, Any]:
    next_payload = dict(payload)
    metadata = next_payload.get("metadata")
    if not isinstance(metadata, dict):
        metadata = {}
    metadata.update(
        {
            "renderer": "markdown",
            "preview_format": "MARKDOWN",
            "render_kernel": "react-markdown",
        }
    )
    next_payload["metadata"] = metadata
    next_payload["markdown"] = "\n".join(_markdown_lines("DOCUMENT", title, next_payload, citations)).strip() + "\n"
    return next_payload


def _render_markdown(resource_id: str, resource_type: str, title: str, payload: dict[str, Any], citations: list[dict[str, Any]]) -> tuple[str, str]:
    settings = get_settings()
    storage_dir = Path(settings.resource_storage_dir) / "generated" / resource_type.lower()
    storage_dir.mkdir(parents=True, exist_ok=True)
    path = storage_dir / f"{resource_id}.md"
    if resource_type == "DOCUMENT" and isinstance(payload.get("markdown"), str) and payload["markdown"].strip():
        path.write_text(str(payload["markdown"]).strip() + "\n", encoding="utf-8")
    else:
        path.write_text("\n".join(_markdown_lines(resource_type, title, payload, citations)), encoding="utf-8")
    return str(path), "MD"


def _render_generic_resource(resource_id: str, resource_type: str, title: str, payload: dict[str, Any], citations: list[dict[str, Any]]) -> tuple[str, str]:
    if resource_type == "DOCUMENT":
        try:
            return _render_docx(resource_id, title, payload, citations)
        except ApiError as exc:
            detail = exc.detail if isinstance(exc.detail, dict) else {}
            if detail.get("code") == "DOC_RENDERER_NOT_INSTALLED":
                return _render_markdown(resource_id, resource_type, title, payload, citations)
            raise
    return _render_markdown(resource_id, resource_type, title, payload, citations)


def resource_media_type(resource: StudentGeneratedResource) -> str:
    format_map = {
        "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "MD": "text/markdown; charset=utf-8",
        "JSON": "application/json",
    }
    return format_map.get((resource.file_format or "").upper(), "application/octet-stream")


def resource_preview_path(resource: StudentGeneratedResource) -> str | None:
    render_payload = _json_loads(resource.render_payload_json, {})
    if not isinstance(render_payload, dict):
        return None
    metadata = render_payload.get("metadata")
    if not isinstance(metadata, dict):
        return None
    raw_path = metadata.get("preview_path")
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None
    path = Path(raw_path).resolve()
    preview_root = (Path(get_settings().resource_storage_dir) / "generated" / "previews").resolve()
    try:
        path.relative_to(preview_root)
    except ValueError:
        return None
    if not path.exists() or path.suffix.lower() != ".pdf":
        return None
    return str(path)


def ensure_resource_preview(resource: StudentGeneratedResource) -> str | None:
    existing = resource_preview_path(resource)
    if existing:
        return existing
    if resource.resource_type != "PPT" or not resource.file_path:
        return None

    preview_metadata = _render_ppt_pdf_preview(resource.id, resource.file_path, resource.file_format)
    render_payload = _json_loads(resource.render_payload_json, {})
    if not isinstance(render_payload, dict):
        render_payload = {}
    metadata = render_payload.get("metadata")
    if not isinstance(metadata, dict):
        metadata = {}
    metadata.update(preview_metadata)
    render_payload["metadata"] = metadata
    resource.render_payload_json = _json_dumps(render_payload)
    resource.updated_at = utc_now()
    return resource_preview_path(resource)


def _serialize_resource(resource: StudentGeneratedResource) -> dict[str, Any]:
    render_payload = _json_loads(resource.render_payload_json, {})
    citations = _json_loads(resource.citations_json, [])
    if not isinstance(render_payload, dict):
        render_payload = {}
    metadata = render_payload.get("metadata")
    if isinstance(metadata, dict) and metadata.get("renderer") == "presenton":
        settings = get_settings()
        public_base_url = settings.presenton_public_base_url or settings.presenton_base_url
        edit_path = metadata.get("presenton_edit_path")
        download_path = metadata.get("presenton_download_path")
        if public_base_url and isinstance(edit_path, str) and edit_path and not metadata.get("presenton_edit_url"):
            metadata["presenton_edit_url"] = urljoin(str(public_base_url).rstrip("/") + "/", edit_path.lstrip("/"))
        if public_base_url and isinstance(download_path, str) and download_path and not metadata.get("presenton_download_url"):
            metadata["presenton_download_url"] = urljoin(str(public_base_url).rstrip("/") + "/", download_path.lstrip("/"))
        render_payload["metadata"] = metadata
        if not render_payload.get("presenton_slides") and isinstance(metadata.get("presenton_presentation_id"), str):
            slides_from_presenton = fetch_presenton_slides_sync(str(metadata["presenton_presentation_id"]), settings)
            if slides_from_presenton:
                render_payload["presenton_slides"] = slides_from_presenton
    if isinstance(metadata, dict):
        preview_path = resource_preview_path(resource)
        if preview_path:
            metadata["preview_available"] = True
            metadata["preview_url"] = f"/api/v1/student/resources/{resource.id}/preview"
            metadata["preview_format"] = "PDF"
            render_payload["metadata"] = metadata
    preview_path = resource_preview_path(resource)
    slides = render_payload.get("slides", [])
    sections = render_payload.get("sections", [])
    nodes = render_payload.get("nodes", [])
    questions = render_payload.get("questions", [])
    cards = render_payload.get("cards", [])
    segments = render_payload.get("segments", [])
    item_count = 0
    for collection in (slides, sections, nodes, questions, cards, segments):
        if isinstance(collection, list) and collection:
            item_count = len(collection)
            break
    return {
        "id": resource.id,
        "resource_type": resource.resource_type,
        "resource_type_label": RESOURCE_TYPE_LABELS.get(resource.resource_type, resource.resource_type),
        "title": resource.title,
        "status": resource.status,
        "summary": resource.summary,
        "knowledge_point": resource.knowledge_point,
        "course_id": resource.course_id,
        "run_id": resource.run_id,
        "confidence": resource.confidence,
        "citations": citations,
        "render_payload": render_payload,
        "file_format": resource.file_format,
        "slide_count": len(slides) if isinstance(slides, list) else 0,
        "item_count": item_count,
        "download_available": bool(resource.file_path),
        "preview_available": bool(preview_path),
        "preview_url": f"/api/v1/student/resources/{resource.id}/preview" if preview_path else None,
        "saved_to_resource_center": resource.saved_to_resource_center,
        "created_at": iso(resource.created_at),
        "updated_at": iso(resource.updated_at),
        "saved_at": iso(resource.saved_at),
    }


def _create_run_node(db: Session, state: PptResourceState) -> PptResourceState:
    context = AgentRunContext(
        run_id=state["run_id"],
        workflow_type=WORKFLOW_TYPE,
        student_id=state["student_id"],
        course_id=state["course_id"],
    )
    run = start_run(
        db,
        context,
        input_payload={
            "resource_type": "PPT",
            "message": _trim(state["message"], 300),
            "session_id": state.get("session_id"),
        },
        model_provider="OPENAI_COMPATIBLE" if get_settings().model_api_key else "RULE_FALLBACK",
        model_name=get_settings().model_name or "rule-template",
        prompt_version=PROMPT_VERSION,
    )
    state["run"] = run
    record_step(db, run, step_name="create_run", step_order=1, output_summary={"run_id": run.id})
    return state


def _context_node(db: Session, user: User, course: Course, state: PptResourceState) -> PptResourceState:
    profile = serialize_learner_profile(
        db,
        student_id=user.id,
        course_id=course.id,
        class_id=state["class_id"],
    )
    sources = _load_sources(db, course.id, state["message"])
    knowledge_point = _guess_knowledge_point(state["message"], sources)
    state["profile"] = profile
    state["sources"] = sources
    state["knowledge_point"] = knowledge_point
    record_step(
        db,
        state["run"],
        step_name="build_context",
        step_order=2,
        output_summary={
            "profile_available": profile is not None,
            "source_ids": [source.id for source in sources],
            "knowledge_point": knowledge_point,
        },
    )
    return state


async def _content_node(db: Session, course: Course, state: PptResourceState) -> PptResourceState:
    sources = state.get("sources", [])
    try:
        model_result = await _generate_model_slides(
            course=course,
            message=state["message"],
            knowledge_point=state["knowledge_point"],
            profile=state.get("profile"),
            sources=sources,
        )
    except LLMError as exc:
        model_result = None
        reason = exc.detail or str(exc)
        state["model_content_fallback_error"] = reason
        record_step(
            db,
            state["run"],
            step_name="model_content_generation",
            step_order=3,
            status="FAILED",
            output_summary={"fallback": True, "reason": reason[:240], "llm_error_code": exc.code},
        )
    except Exception as exc:
        model_result = None
        state["model_content_fallback_error"] = str(exc)
        record_step(
            db,
            state["run"],
            step_name="model_content_generation",
            step_order=3,
            status="FAILED",
            output_summary={"fallback": True, "reason": str(exc)[:240]},
        )
    else:
        record_step(
            db,
            state["run"],
            step_name="model_content_generation",
            step_order=3,
            output_summary={"used_model": model_result is not None},
        )

    if model_result is None:
        model_result = _fallback_slides(
            message=state["message"],
            course=course,
            knowledge_point=state["knowledge_point"],
            sources=sources,
            profile=state.get("profile"),
        )
    title, summary, slides = model_result
    citations = [_citation(source) for source in sources]
    state["title"] = title
    state["summary"] = summary
    state["slides"] = slides
    state["citations"] = citations
    state["confidence"] = 0.86 if citations else 0.52
    record_step(
        db,
        state["run"],
        step_name="validate_content",
        step_order=4,
        output_summary={"slide_count": len(slides), "citation_count": len(citations)},
    )
    return state


async def _render_node(db: Session, state: PptResourceState) -> PptResourceState:
    resource_id = _new_id("res")
    path, file_format, render_metadata = await _render_ppt_resource_file(db, state, resource_id)
    if state.get("model_content_fallback_error"):
        render_metadata.update(
            {
                "model_content_fallback": True,
                "model_content_fallback_error": state["model_content_fallback_error"],
            }
        )
    preview_metadata = _render_ppt_pdf_preview(resource_id, path, file_format)
    render_metadata.update(preview_metadata)
    record_step(
        db,
        state["run"],
        step_name="render_ppt_pdf_preview",
        step_order=6,
        status="SUCCEEDED" if preview_metadata.get("preview_available") else "FAILED",
        output_summary={
            "resource_id": resource_id,
            "preview_format": preview_metadata.get("preview_format"),
            "preview_available": bool(preview_metadata.get("preview_available")),
            "fallback": not bool(preview_metadata.get("preview_available")),
            "reason": preview_metadata.get("preview_error"),
        },
    )
    resource = StudentGeneratedResource(
        id=resource_id,
        student_id=state["student_id"],
        course_id=state["course_id"],
        class_id=state["class_id"],
        run_id=state["run_id"],
        session_id=state.get("session_id"),
        resource_type="PPT",
        title=state["title"],
        prompt=state["message"],
        knowledge_point=state["knowledge_point"],
        summary=state["summary"],
        status="READY",
        render_payload_json=_json_dumps(
            {
                "slides": state["slides"],
                "presenton_slides": state.get("presenton_slides", []),
                "metadata": render_metadata,
            }
        ),
        citations_json=_json_dumps(state["citations"]),
        file_path=path,
        file_format=file_format,
        confidence=state["confidence"],
        saved_to_resource_center=False,
    )
    db.add(resource)
    state["resource"] = resource
    state["file_path"] = path
    record_step(
        db,
        state["run"],
        step_name="persist_ppt_resource",
        step_order=8,
        output_summary={"resource_id": resource_id, "file_format": file_format, "renderer": render_metadata.get("renderer")},
    )
    return state


def _validate_resource_type(resource_type: str) -> str:
    normalized = (resource_type or "").strip().upper()
    if normalized not in SUPPORTED_RESOURCE_TYPES:
        raise ApiError(
            400,
            "UNSUPPORTED_RESOURCE_TYPE",
            "暂不支持该资源类型。",
            details={"resource_type": resource_type, "supported": sorted(SUPPORTED_RESOURCE_TYPES)},
        )
    return normalized


def _create_generic_run_node(db: Session, state: GenericResourceState) -> GenericResourceState:
    context = AgentRunContext(
        run_id=state["run_id"],
        workflow_type=GENERIC_WORKFLOW_TYPE,
        student_id=state["student_id"],
        course_id=state["course_id"],
    )
    run = start_run(
        db,
        context,
        input_payload={
            "resource_type": state["resource_type"],
            "message": _trim(state["message"], 300),
            "session_id": state.get("session_id"),
        },
        model_provider="RULE_FALLBACK",
        model_name="rule-template",
        prompt_version=GENERIC_PROMPT_VERSION,
    )
    state["run"] = run
    record_step(db, run, step_name="create_run", step_order=1, output_summary={"run_id": run.id})
    return state


def _generic_context_node(db: Session, user: User, course: Course, state: GenericResourceState) -> GenericResourceState:
    profile = serialize_learner_profile(
        db,
        student_id=user.id,
        course_id=course.id,
        class_id=state["class_id"],
    )
    sources = _load_sources(db, course.id, state["message"])
    knowledge_point = _guess_knowledge_point(state["message"], sources)
    state["profile"] = profile
    state["sources"] = sources
    state["knowledge_point"] = knowledge_point
    record_step(
        db,
        state["run"],
        step_name="build_context",
        step_order=2,
        output_summary={
            "profile_available": profile is not None,
            "source_ids": [source.id for source in sources],
            "knowledge_point": knowledge_point,
        },
    )
    return state


def _generic_content_node(db: Session, course: Course, state: GenericResourceState) -> GenericResourceState:
    sources = state.get("sources", [])
    title, summary, render_payload, item_count = _fallback_resource_payload(
        resource_type=state["resource_type"],
        message=state["message"],
        course=course,
        knowledge_point=state["knowledge_point"],
        sources=sources,
    )
    citations = [_citation(source) for source in sources]
    if state["resource_type"] == "DOCUMENT":
        render_payload = _attach_document_markdown_preview(title, render_payload, citations)
    state["title"] = title
    state["summary"] = summary
    state["render_payload"] = render_payload
    state["citations"] = citations
    state["item_count"] = item_count
    state["confidence"] = 0.82 if citations else 0.5
    record_step(
        db,
        state["run"],
        step_name="generate_structured_content",
        step_order=3,
        output_summary={
            "resource_type": state["resource_type"],
            "item_count": item_count,
            "citation_count": len(citations),
        },
    )
    return state


def _generic_render_node(db: Session, state: GenericResourceState) -> GenericResourceState:
    resource_id = _new_id("res")
    path, file_format = _render_generic_resource(
        resource_id,
        state["resource_type"],
        state["title"],
        state["render_payload"],
        state["citations"],
    )
    resource = StudentGeneratedResource(
        id=resource_id,
        student_id=state["student_id"],
        course_id=state["course_id"],
        class_id=state["class_id"],
        run_id=state["run_id"],
        session_id=state.get("session_id"),
        resource_type=state["resource_type"],
        title=state["title"],
        prompt=state["message"],
        knowledge_point=state["knowledge_point"],
        summary=state["summary"],
        status="READY",
        render_payload_json=_json_dumps(state["render_payload"]),
        citations_json=_json_dumps(state["citations"]),
        file_path=path,
        file_format=file_format,
        confidence=state["confidence"],
        saved_to_resource_center=False,
    )
    db.add(resource)
    state["resource"] = resource
    state["file_path"] = path
    state["file_format"] = file_format
    record_step(
        db,
        state["run"],
        step_name="render_resource_file_and_preview",
        step_order=4,
        output_summary={"resource_id": resource_id, "file_format": file_format},
    )
    return state


def _create_mind_map_run_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    context = AgentRunContext(
        run_id=state["run_id"],
        workflow_type=MIND_MAP_WORKFLOW_TYPE,
        student_id=state["student_id"],
        course_id=state["course_id"],
    )
    run = start_run(
        db,
        context,
        input_payload={
            "resource_type": "MIND_MAP",
            "message": _trim(state["message"], 300),
            "session_id": state.get("session_id"),
        },
        model_provider="RULE_FALLBACK",
        model_name="mind-map-agent-template",
        prompt_version=MIND_MAP_PROMPT_VERSION,
    )
    state["run"] = run
    record_step(db, run, step_name="create_run", step_order=1, output_summary={"run_id": run.id})
    return state


def _mind_map_context_node(db: Session, user: User, course: Course, state: MindMapResourceState) -> MindMapResourceState:
    profile = serialize_learner_profile(
        db,
        student_id=user.id,
        course_id=course.id,
        class_id=state["class_id"],
    )
    sources = _load_sources(db, course.id, state["message"])
    knowledge_point = _guess_knowledge_point(state["message"], sources)
    state["profile"] = profile
    state["profile_focus"] = _profile_focus(profile)
    state["sources"] = sources
    state["knowledge_point"] = knowledge_point
    record_step(
        db,
        state["run"],
        step_name="build_context",
        step_order=2,
        output_summary={
            "profile_available": profile is not None,
            "source_ids": [source.id for source in sources],
            "knowledge_point": knowledge_point,
        },
    )
    return state


def _mind_map_planner_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    focus = state.get("profile_focus") or {"weak_points": [], "frequent_errors": [], "recommendations": []}
    topic_family = _mind_map_topic_family(state["knowledge_point"])
    top_level_branches = (
        ["语言基础", "流程控制", "函数与对象", "常见错误", "画像提醒", "练习路径"]
        if topic_family == "programming"
        else ["核心概念", "建模流程", "常见误区", "画像提醒", "练习路径"]
        if topic_family == "machine_learning"
        else ["核心概念", "关键过程", "边界条件", "画像提醒", "练习路径"]
    )
    plan = {
        "central_topic": f"{state['knowledge_point']}学习地图",
        "map_purpose": "帮助学生梳理概念、过程、易错点、画像提醒和练习路径。",
        "target_depth": 2,
        "top_level_branches": top_level_branches,
        "required_relationship_types": ["contains", "causes", "solves", "example_of", "common_mistake", "next_step"],
        "coverage_requirements": [
            "必须回应用户原始生成要求",
            "至少包含一个易错点分支",
            "至少包含一个学习画像提醒分支",
            "节点必须适合前端预览和资源中心导出",
        ],
        "profile_focus": focus,
        "topic_family": topic_family,
    }
    state["plan"] = plan
    record_step(
        db,
        state["run"],
        step_name="mind_map_planner_agent",
        step_order=3,
        output_summary={
            "central_topic": plan["central_topic"],
            "branch_count": len(plan["top_level_branches"]),
            "profile_focus_used": bool(focus.get("weak_points") or focus.get("frequent_errors")),
        },
    )
    return state


def _mind_map_content_node(db: Session, course: Course, state: MindMapResourceState) -> MindMapResourceState:
    title, summary, render_payload, item_count = _fallback_mind_map(
        message=state["message"],
        course=course,
        knowledge_point=state["knowledge_point"],
        sources=state.get("sources", []),
        profile=state.get("profile"),
    )
    state["title"] = title
    state["summary"] = summary
    state["render_payload"] = render_payload
    state["item_count"] = item_count
    record_step(
        db,
        state["run"],
        step_name="mind_map_content_agent",
        step_order=4,
        output_summary={"node_count": item_count, "edge_count": len(render_payload.get("edges", []))},
    )
    return state


def _mind_map_relationship_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    payload = state["render_payload"]
    metadata = payload.get("metadata") if isinstance(payload.get("metadata"), dict) else {}
    topic_family = str(metadata.get("topic_family") or _mind_map_topic_family(state["knowledge_point"]))
    edges = payload.get("edges", [])
    if not isinstance(edges, list):
        edges = []
    existing_pairs = {
        (str(edge.get("source") or edge.get("source_node_id")), str(edge.get("target") or edge.get("target_node_id")))
        for edge in edges
        if isinstance(edge, dict)
    }
    if topic_family == "programming":
        extra_edges = [
            ("concept_1", "procedure_1", "prerequisite", "变量类型支撑分支"),
            ("procedure_2", "mistake_3", "causes", "循环结构暴露符号和边界问题"),
            ("mistake_1", "practice_2", "solves", "用改错题巩固"),
        ]
    elif topic_family == "machine_learning":
        extra_edges = [
            ("concept_1", "procedure_1", "prerequisite", "先理解概念再看流程"),
            ("procedure_2", "mistake_1", "causes", "评估表现暴露泛化问题"),
            ("mistake_2", "practice_1", "solves", "用辨析题巩固"),
        ]
    else:
        extra_edges = [
            ("concept_1", "procedure_1", "prerequisite", "先理解再操作"),
            ("procedure_2", "boundary_3", "causes", "连续状态影响边界"),
            ("boundary_2", "practice_2", "solves", "用伪过程检查"),
        ]
    for source_id, target_id, relation, label in extra_edges:
        if (source_id, target_id) in existing_pairs:
            continue
        edges.append(
            {
                "source": source_id,
                "target": target_id,
                "source_node_id": source_id,
                "target_node_id": target_id,
                "relationship_type": relation,
                "label": label,
            }
        )
    payload["edges"] = edges
    state["render_payload"] = payload
    record_step(
        db,
        state["run"],
        step_name="relationship_refiner_agent",
        step_order=5,
        output_summary={"edge_count": len(edges), "cross_relationships": 3},
    )
    return state


def _mind_map_guard_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    allowed_source_ids = {source.id for source in state.get("sources", [])}
    payload = state["render_payload"]
    invalid_ids: set[str] = set()
    for node in payload.get("nodes", []) if isinstance(payload.get("nodes"), list) else []:
        if not isinstance(node, dict):
            continue
        raw_ids = node.get("citation_ids") if isinstance(node.get("citation_ids"), list) else node.get("citations")
        valid_ids = [str(source_id) for source_id in (raw_ids or []) if str(source_id) in allowed_source_ids]
        invalid_ids.update(str(source_id) for source_id in (raw_ids or []) if str(source_id) not in allowed_source_ids)
        node["citation_ids"] = valid_ids
        node["citations"] = valid_ids
        if not valid_ids and node.get("node_type") in {"concept", "procedure", "mistake"}:
            node["confidence"] = min(float(node.get("confidence") or 0.5), 0.58)
    citations = [_citation(source) for source in state.get("sources", [])]
    risk_flags = list(payload.get("risk_flags", [])) if isinstance(payload.get("risk_flags"), list) else []
    if not citations and "引用不足" not in risk_flags:
        risk_flags.append("引用不足")
    if invalid_ids and "存在无效引用并已过滤" not in risk_flags:
        risk_flags.append("存在无效引用并已过滤")
    payload["risk_flags"] = risk_flags
    state["render_payload"] = payload
    state["citations"] = citations
    state["risk_flags"] = risk_flags
    record_step(
        db,
        state["run"],
        step_name="citation_guard_agent",
        step_order=6,
        output_summary={
            "citation_count": len(citations),
            "invalid_citation_ids": sorted(invalid_ids),
            "risk_flags": risk_flags,
        },
    )
    return state


def _mind_map_critic_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    payload = state["render_payload"]
    nodes = payload.get("nodes", []) if isinstance(payload.get("nodes"), list) else []
    level_one_count = sum(1 for node in nodes if isinstance(node, dict) and int(node.get("level", 0) or 0) == 1)
    has_mistake = any(isinstance(node, dict) and node.get("node_type") == "mistake" for node in nodes)
    has_profile_tip = any(isinstance(node, dict) and node.get("node_type") == "profile_tip" for node in nodes)
    risk_flags = list(state.get("risk_flags", []))
    if not 4 <= level_one_count <= 7 and "一级分支数量异常" not in risk_flags:
        risk_flags.append("一级分支数量异常")
    if not has_mistake and "缺少易错点" not in risk_flags:
        risk_flags.append("缺少易错点")
    if not has_profile_tip and "缺少画像提醒" not in risk_flags:
        risk_flags.append("缺少画像提醒")
    payload["risk_flags"] = risk_flags
    payload.setdefault("metadata", {})
    if isinstance(payload["metadata"], dict):
        payload["metadata"].update(
            {
                "critic": "rule_fallback_mind_map_critic",
                "quality_checks": {
                    "node_count": len(nodes),
                    "level_one_count": level_one_count,
                    "has_mistake": has_mistake,
                    "has_profile_tip": has_profile_tip,
                },
            }
        )
    state["render_payload"] = payload
    state["risk_flags"] = risk_flags
    state["confidence"] = 0.84 if state.get("citations") and not risk_flags else 0.58 if state.get("citations") else 0.5
    record_step(
        db,
        state["run"],
        step_name="mind_map_critic_agent",
        step_order=7,
        status="SUCCEEDED",
        output_summary={
            "node_count": len(nodes),
            "level_one_count": level_one_count,
            "risk_flags": risk_flags,
            "confidence": state["confidence"],
        },
    )
    return state


def _mind_map_persist_node(db: Session, state: MindMapResourceState) -> MindMapResourceState:
    resource_id = _new_id("res")
    path, file_format = _render_generic_resource(
        resource_id,
        "MIND_MAP",
        state["title"],
        state["render_payload"],
        state.get("citations", []),
    )
    resource = StudentGeneratedResource(
        id=resource_id,
        student_id=state["student_id"],
        course_id=state["course_id"],
        class_id=state["class_id"],
        run_id=state["run_id"],
        session_id=state.get("session_id"),
        resource_type="MIND_MAP",
        title=state["title"],
        prompt=state["message"],
        knowledge_point=state["knowledge_point"],
        summary=state["summary"],
        status="READY",
        render_payload_json=_json_dumps(state["render_payload"]),
        citations_json=_json_dumps(state.get("citations", [])),
        file_path=path,
        file_format=file_format,
        confidence=state["confidence"],
        saved_to_resource_center=False,
    )
    db.add(resource)
    state["resource"] = resource
    state["file_path"] = path
    state["file_format"] = file_format
    record_step(
        db,
        state["run"],
        step_name="persist_mind_map_resource",
        step_order=8,
        output_summary={"resource_id": resource_id, "file_format": file_format},
    )
    return state


async def generate_ppt_resource(
    db: Session,
    *,
    user: User,
    class_id: str,
    course: Course,
    message: str,
    session_id: str | None = None,
) -> dict[str, Any]:
    state: PptResourceState = {
        "run_id": new_run_id(),
        "student_id": user.id,
        "class_id": class_id,
        "course_id": course.id,
        "session_id": session_id,
        "message": message.strip(),
    }

    if StateGraph is not None:
        # Keep the graph boundary explicit for future expansion; DB-bound nodes still
        # execute through local closures so permissions and transactions stay in FastAPI.
        async def generate_content_node(graph_state: PptResourceState) -> PptResourceState:
            return await _content_node(db, course, graph_state)

        async def render_resource_node(graph_state: PptResourceState) -> PptResourceState:
            return await _render_node(db, graph_state)

        graph = StateGraph(PptResourceState)
        graph.add_node("create_run", lambda graph_state: _create_run_node(db, graph_state))
        graph.add_node("build_context", lambda graph_state: _context_node(db, user, course, graph_state))
        graph.add_node("generate_content", generate_content_node)
        graph.add_node("render_resource", render_resource_node)
        graph.set_entry_point("create_run")
        graph.add_edge("create_run", "build_context")
        graph.add_edge("build_context", "generate_content")
        graph.add_edge("generate_content", "render_resource")
        graph.add_edge("render_resource", END)
        state = await graph.compile().ainvoke(state)
    else:
        state = _create_run_node(db, state)
        state = _context_node(db, user, course, state)
        state = await _content_node(db, course, state)
        state = await _render_node(db, state)

    resource = state["resource"]
    finish_run(
        db,
        state["run"],
        output={"resource_id": resource.id, "title": resource.title, "slide_count": len(state["slides"])},
        model_provider="OPENAI_COMPATIBLE" if get_settings().model_api_key else "RULE_FALLBACK",
        model_name=get_settings().model_name or "rule-template",
        prompt_version=PROMPT_VERSION,
    )
    return _serialize_resource(resource)


async def generate_mind_map_resource(
    db: Session,
    *,
    user: User,
    class_id: str,
    course: Course,
    message: str,
    session_id: str | None = None,
) -> dict[str, Any]:
    state: MindMapResourceState = {
        "run_id": new_run_id(),
        "student_id": user.id,
        "class_id": class_id,
        "course_id": course.id,
        "session_id": session_id,
        "message": message.strip(),
        "resource_type": "MIND_MAP",
    }

    if StateGraph is not None:
        graph = StateGraph(MindMapResourceState)
        graph.add_node("create_run", lambda graph_state: _create_mind_map_run_node(db, graph_state))
        graph.add_node("build_context", lambda graph_state: _mind_map_context_node(db, user, course, graph_state))
        graph.add_node("plan_structure", lambda graph_state: _mind_map_planner_node(db, graph_state))
        graph.add_node("generate_nodes", lambda graph_state: _mind_map_content_node(db, course, graph_state))
        graph.add_node("refine_relationships", lambda graph_state: _mind_map_relationship_node(db, graph_state))
        graph.add_node("guard_citations", lambda graph_state: _mind_map_guard_node(db, graph_state))
        graph.add_node("review_quality", lambda graph_state: _mind_map_critic_node(db, graph_state))
        graph.add_node("persist_resource", lambda graph_state: _mind_map_persist_node(db, graph_state))
        graph.set_entry_point("create_run")
        graph.add_edge("create_run", "build_context")
        graph.add_edge("build_context", "plan_structure")
        graph.add_edge("plan_structure", "generate_nodes")
        graph.add_edge("generate_nodes", "refine_relationships")
        graph.add_edge("refine_relationships", "guard_citations")
        graph.add_edge("guard_citations", "review_quality")
        graph.add_edge("review_quality", "persist_resource")
        graph.add_edge("persist_resource", END)
        state = await graph.compile().ainvoke(state)
    else:
        state = _create_mind_map_run_node(db, state)
        state = _mind_map_context_node(db, user, course, state)
        state = _mind_map_planner_node(db, state)
        state = _mind_map_content_node(db, course, state)
        state = _mind_map_relationship_node(db, state)
        state = _mind_map_guard_node(db, state)
        state = _mind_map_critic_node(db, state)
        state = _mind_map_persist_node(db, state)

    resource = state["resource"]
    finish_run(
        db,
        state["run"],
        output={
            "resource_id": resource.id,
            "title": resource.title,
            "resource_type": resource.resource_type,
            "node_count": state["item_count"],
            "risk_flags": state.get("risk_flags", []),
        },
        model_provider="RULE_FALLBACK",
        model_name="mind-map-agent-template",
        prompt_version=MIND_MAP_PROMPT_VERSION,
    )
    return _serialize_resource(resource)


async def generate_learning_resource(
    db: Session,
    *,
    user: User,
    class_id: str,
    course: Course,
    message: str,
    resource_type: str,
    session_id: str | None = None,
) -> dict[str, Any]:
    normalized_type = _validate_resource_type(resource_type)
    if normalized_type == "PPT":
        return await generate_ppt_resource(
            db,
            user=user,
            class_id=class_id,
            course=course,
            message=message,
            session_id=session_id,
        )
    if normalized_type == "MIND_MAP":
        return await generate_mind_map_resource(
            db,
            user=user,
            class_id=class_id,
            course=course,
            message=message,
            session_id=session_id,
        )

    state: GenericResourceState = {
        "run_id": new_run_id(),
        "student_id": user.id,
        "class_id": class_id,
        "course_id": course.id,
        "session_id": session_id,
        "message": message.strip(),
        "resource_type": normalized_type,
    }

    if StateGraph is not None:
        graph = StateGraph(GenericResourceState)
        graph.add_node("create_run", lambda graph_state: _create_generic_run_node(db, graph_state))
        graph.add_node("build_context", lambda graph_state: _generic_context_node(db, user, course, graph_state))
        graph.add_node("generate_content", lambda graph_state: _generic_content_node(db, course, graph_state))
        graph.add_node("render_resource", lambda graph_state: _generic_render_node(db, graph_state))
        graph.set_entry_point("create_run")
        graph.add_edge("create_run", "build_context")
        graph.add_edge("build_context", "generate_content")
        graph.add_edge("generate_content", "render_resource")
        graph.add_edge("render_resource", END)
        state = await graph.compile().ainvoke(state)
    else:
        state = _create_generic_run_node(db, state)
        state = _generic_context_node(db, user, course, state)
        state = _generic_content_node(db, course, state)
        state = _generic_render_node(db, state)

    resource = state["resource"]
    finish_run(
        db,
        state["run"],
        output={
            "resource_id": resource.id,
            "title": resource.title,
            "resource_type": resource.resource_type,
            "item_count": state["item_count"],
        },
        model_provider="RULE_FALLBACK",
        model_name="rule-template",
        prompt_version=GENERIC_PROMPT_VERSION,
    )
    return _serialize_resource(resource)


def get_generated_resource(db: Session, *, student_id: str, resource_id: str) -> StudentGeneratedResource:
    resource = db.get(StudentGeneratedResource, resource_id)
    if resource is None or resource.student_id != student_id:
        raise ApiError(404, "GENERATED_RESOURCE_NOT_FOUND", "生成资源不存在或不可访问。")
    return resource


def save_generated_resource(db: Session, *, user: User, class_id: str, resource_id: str) -> dict[str, Any]:
    resource = get_generated_resource(db, student_id=user.id, resource_id=resource_id)
    if not resource.saved_to_resource_center:
        resource.saved_to_resource_center = True
        resource.saved_at = utc_now()
        resource.updated_at = utc_now()
        db.add(
            LearnerEvent(
                id=_new_id("event"),
                student_id=user.id,
                course_id=resource.course_id,
                class_id=class_id,
                event_type="artifact_saved",
                knowledge_points=_json_dumps([resource.knowledge_point] if resource.knowledge_point else []),
                payload=_json_dumps(
                    {
                        "resource_id": resource.id,
                        "resource_type": resource.resource_type,
                        "source": "ai_resource_card",
                    }
                ),
            )
        )
    return _serialize_resource(resource)


def list_saved_generated_resources(
    db: Session,
    *,
    student_id: str,
    course_id: str | None = None,
    limit: int = 50,
) -> list[dict[str, Any]]:
    statement = select(StudentGeneratedResource).where(
        StudentGeneratedResource.student_id == student_id,
        StudentGeneratedResource.saved_to_resource_center.is_(True),
    )
    if course_id:
        statement = statement.where(StudentGeneratedResource.course_id == course_id)
    resources = list(
        db.scalars(
            statement.order_by(StudentGeneratedResource.saved_at.desc(), StudentGeneratedResource.updated_at.desc()).limit(limit)
        ).all()
    )
    return [_serialize_resource(resource) for resource in resources]
